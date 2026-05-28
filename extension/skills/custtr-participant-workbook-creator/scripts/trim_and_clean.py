"""Remove the last 2 slides from every .pptx in a folder, except the last file.

Usage:
    uv run --with python-pptx python trim_slides.py "<folder_path>"

Files are sorted alphabetically. The last file (alphabetically) is left untouched.
All others have their final 2 slides removed and are saved in-place.
"""
import sys
import os
from pathlib import Path
from pptx import Presentation

sys.stdout.reconfigure(encoding='utf-8')

if len(sys.argv) < 2:
    print("Usage: trim_slides.py <folder_path>")
    sys.exit(1)

folder = Path(sys.argv[1])
files = sorted(folder.glob("*.pptx"))

if not files:
    print("No .pptx files found.")
    sys.exit(1)

print(f"Found {len(files)} files. Last file (kept intact): {files[-1].name}\n")

results = []
for i, f in enumerate(files):
    prs = Presentation(str(f))
    original_count = len(prs.slides)
    if i < len(files) - 1:
        # Remove last 2 slides
        xml_slides = prs.slides._sldIdLst
        to_remove = list(xml_slides)[-2:]
        for sldId in to_remove:
            rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
            prs.part.drop_rel(rId)
            xml_slides.remove(sldId)
        prs.save(str(f))
        new_count = original_count - 2
        print(f"  TRIMMED  {f.name}: {original_count} -> {new_count} slides")
    else:
        new_count = original_count
        print(f"  KEPT     {f.name}: {original_count} slides (last file, untouched)")
    results.append((f.name, new_count))

print(f"\nSlide counts after trim: {[c for _, c in results]}")
print(f"Total slides: {sum(c for _, c in results)}")
