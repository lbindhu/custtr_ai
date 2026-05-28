"""
SB Cleanup Script — applies 13 cleanup rules to a PowerPoint SB file.

Rules:
 1. Remove all comments from slides
 2. Slide after LO = Slide-3; LO itself = Slide-2; same title/click slides share same number
 3. Remove all strikethrough runs + collapse resulting double spaces
 4. Remove highlight formatting; fix font to black/white for readability on white bg
 5. Remove empty paragraphs left after strikeout removal (realign)
 6+9. Remove ALL shapes outside the slide boundary, except yellow NTD/Shared boxes
 7. Remove complete slides that have a diagonal deletion line OR all-strikethrough text
 8. Remove double spaces between words in slide text and notes
10. Remove yellow boxes with strikethrough AND red-colored text
11. Do NOT remove yellow boxes containing NTD / Fully Shared / Partially Shared text
12. Output filename = input name with version token (_V<n>) stripped
13. Remove yellow boxes containing animation notes (authoring cues for developers)
14. Remove yellow boxes that are empty after strikethrough removal
"""

import re
import sys
from pptx import Presentation
from pptx.oxml.ns import qn as _qn
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

# ── XML Namespaces ─────────────────────────────────────────────────────────────
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

# ── Patterns ───────────────────────────────────────────────────────────────────
SLIDE_NUM_PATTERN = re.compile(
    r"^(Slide[-\s]\d+(\s*([\n\r\x0b/\\]\s*)?Branch\s*[\d.]+)?|"
    r"New\s*[\n\r\x0b]?\s*Slide\s+\d+|Slide\s+\d+)$",
    re.IGNORECASE | re.DOTALL,
)
LO_PATTERN  = re.compile(r"learning\s+obj|^objectives?$", re.IGNORECASE)
AYK_PATTERN = re.compile(r"apply\s+your\s+knowledge", re.IGNORECASE)
# Rule 11: protect these yellow boxes from Rule 10 removal
NTD_PATTERN = re.compile(
    r"(NTD\s*:|Fully\s+Shared|Partially\s+Shared|Shared\s+Slide)",
    re.IGNORECASE,
)
# Rule 13: animation notes authoring boxes to remove
ANIMATION_NOTES_PATTERN = re.compile(
    r"^\s*animation\s*(notes?|cues?|instructions?|details?)?[\s:\-]",
    re.IGNORECASE,
)
VERSION_PATTERN = re.compile(r"_V\d+(?=\.pptx$)", re.IGNORECASE)

CLICK_PATTERNS = [
    re.compile(r"click\s+each", re.IGNORECASE),
    re.compile(r"click\s+through", re.IGNORECASE),
    re.compile(r"click\s+the\s+arrow", re.IGNORECASE),
    re.compile(r"\bclick\s+\w+", re.IGNORECASE),
    re.compile(r"drag[- ]and[- ]drop", re.IGNORECASE),
    re.compile(r"\bdrag\b.{0,60}\bdrop\b", re.IGNORECASE | re.DOTALL),
    re.compile(r"match\s+(each|the)", re.IGNORECASE),
    re.compile(r"hover\s+over", re.IGNORECASE),
    re.compile(r"roll\s+over", re.IGNORECASE),
    re.compile(r"tap\s+on", re.IGNORECASE),
]


# ── Color Helpers ──────────────────────────────────────────────────────────────

# AMD PowerPoint theme color → hex mapping (used to resolve schemeClr fills)
AMD_SCHEME_COLORS = {
    "accent1": "E20000",  # AMD red
    "accent2": "282D3F",  # dark navy
    "accent3": "006EB4",  # blue
    "accent4": "00A870",  # green
    "accent5": "0D9079",  # teal
    "accent6": "7B5EA7",  # purple
    "dk1":     "000000",  # black
    "dk2":     "44546A",  # dark grey
    "lt1":     "FFFFFF",  # white
    "lt2":     "E7E6E6",  # light grey
    "bg1":     "FFFFFF",  # white
    "bg2":     "E7E6E6",  # light grey
}


def _lum(hex6):
    """Perceived luminance 0–255."""
    r, g, b = int(hex6[0:2], 16), int(hex6[2:4], 16), int(hex6[4:6], 16)
    return 0.299 * r + 0.587 * g + 0.114 * b


def _is_yellow(h):
    if not h or len(h) != 6:
        return False
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return r > 200 and g > 200 and b < 80


def _is_blue(h):
    if not h or len(h) != 6:
        return False
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    # B dominant; covers #0070C0, #4472C4, #00B0F0, #1F497D, #2E75B6, etc.
    return b > 120 and b > r + 30 and b > g + 20


def _is_red(h):
    if not h or len(h) != 6:
        return False
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return r > 180 and g < 110 and b < 110


def _get_fill_hex(el):
    """
    First solid fill color of this element.
    Returns explicit srgbClr hex, or AMD_SCHEME_COLORS resolution for schemeClr.
    """
    for sf in el.iter(f"{{{NS_A}}}solidFill"):
        for srgb in sf.iter(f"{{{NS_A}}}srgbClr"):
            v = srgb.get("val", "")
            if len(v) == 6:
                return v.upper()
        for sc in sf.iter(f"{{{NS_A}}}schemeClr"):
            val = sc.get("val", "")
            if val in AMD_SCHEME_COLORS:
                return AMD_SCHEME_COLORS[val]
    return None


def _get_line_color_hex(shape_el):
    """Return the stroke/line color of a shape, or None."""
    # Look inside <a:ln> > <a:solidFill> > <a:srgbClr>
    for ln in shape_el.iter(f"{{{NS_A}}}ln"):
        for sf in ln.iter(f"{{{NS_A}}}solidFill"):
            for srgb in sf.iter(f"{{{NS_A}}}srgbClr"):
                v = srgb.get("val", "")
                if len(v) == 6:
                    return v.upper()
    return None


def _get_run_color_hex(run_el):
    rPr = run_el.find(f"{{{NS_A}}}rPr")
    if rPr is None:
        return None
    for sf in rPr.iter(f"{{{NS_A}}}solidFill"):
        for srgb in sf.iter(f"{{{NS_A}}}srgbClr"):
            v = srgb.get("val", "")
            if len(v) == 6:
                return v.upper()
    return None


def _is_strikethrough(run_el):
    rPr = run_el.find(f"{{{NS_A}}}rPr")
    if rPr is None:
        return False
    return rPr.get("strike", "noStrike") in ("sngStrike", "dblStrike")


def _shape_text(shape):
    try:
        if shape.has_text_frame:
            return "\n".join(p.text for p in shape.text_frame.paragraphs).strip()
    except Exception:
        pass
    return ""


def _slide_title(slide):
    try:
        t = slide.shapes.title
        if t and t.has_text_frame:
            return t.text_frame.text.strip()
    except Exception:
        pass
    return ""


# ── Rule 1: Remove Comments ────────────────────────────────────────────────────

def remove_comments(slide, sn):
    removed = 0
    for cm in slide._element.findall(f"{{{NS_P}}}cmLst"):
        slide._element.remove(cm)
        removed += 1
    rels_to_drop = [
        rId for rId, rel in slide.part.rels.items()
        if "comment" in rel.reltype.lower()
    ]
    for rId in rels_to_drop:
        try:
            slide.part.drop_rel(rId)
            removed += 1
        except Exception:
            pass
    if removed:
        print(f"  [R1]  Slide {sn}: {removed} comment(s) removed")
    return removed


# Pattern for "delete/remove this slide" in comments
DELETE_COMMENT_PATTERN = re.compile(
    r"\b(delete|remove)\s+(this\s+)?slide\b",
    re.IGNORECASE,
)


def _slide_has_delete_comment(slide):
    """
    True if any comment on this slide contains text matching
    'delete this slide', 'remove this slide', 'delete slide', 'remove slide', etc.
    Checks both inline <p:cmLst> comments and relationship-based comment parts.
    """
    # 1. Inline comments: <p:cmLst><p:cm><p:text>...</p:text>
    for cm in slide._element.iter(f"{{{NS_P}}}cm"):
        txt_el = cm.find(f"{{{NS_P}}}text")
        if txt_el is not None and txt_el.text:
            if DELETE_COMMENT_PATTERN.search(txt_el.text):
                return True

    # 2. Relationship-based comment parts (modern PPTX comment format)
    for rId, rel in slide.part.rels.items():
        if "comment" not in rel.reltype.lower():
            continue
        try:
            part = rel.target_part
            xml_bytes = part.blob
            # Parse and search for the delete pattern in the raw XML text
            root = etree.fromstring(xml_bytes)
            for el in root.iter():
                if el.text and DELETE_COMMENT_PATTERN.search(el.text):
                    return True
                if el.tail and DELETE_COMMENT_PATTERN.search(el.tail):
                    return True
        except Exception:
            pass
    return False


# ── Rule 7: Detect deleted slides ─────────────────────────────────────────────

def _has_deletion_line(slide, slide_w, slide_h):
    """
    True if the slide has a LINE or CONNECTOR shape that looks like a
    diagonal deletion mark — spanning at least 60% of BOTH the slide width
    AND height simultaneously (i.e. a true diagonal, not a side bracket).

    Color is not checked — deletion lines may be red, teal, or any scheme
    color depending on the authoring team's convention. Only geometry matters.
    A real diagonal deletion line must span most of the slide in both axes.
    """
    NS_P_local = "http://schemas.openxmlformats.org/presentationml/2006/main"

    def _is_diagonal_deletion(w, h):
        """True if the bounding box spans ≥60% of slide in BOTH dimensions."""
        return w >= slide_w * 0.60 and h >= slide_h * 0.60

    # 1. python-pptx LINE shapes (type 9)
    for shape in slide.shapes:
        if shape.shape_type not in (9,):
            continue
        try:
            w = abs(shape.width or 0)
            h = abs(shape.height or 0)
            if _is_diagonal_deletion(w, h):
                return True
        except Exception:
            pass

    # 2. Raw <p:cxnSp> connector elements (sometimes not surfaced as type 9)
    for cxn in slide._element.iter(f"{{{NS_P_local}}}cxnSp"):
        try:
            sp_pr = cxn.find(f"{{{NS_P_local}}}spPr")
            if sp_pr is None:
                sp_pr = cxn.find(f"{{{NS_A}}}spPr")
            if sp_pr is None:
                continue
            xfrm = sp_pr.find(f"{{{NS_A}}}xfrm")
            if xfrm is None:
                continue
            ext = xfrm.find(f"{{{NS_A}}}ext")
            if ext is None:
                continue
            w = abs(int(ext.get("cx", 0)))
            h = abs(int(ext.get("cy", 0)))
            if _is_diagonal_deletion(w, h):
                return True
        except Exception:
            pass
    return False


def _all_text_runs_with_content(slide):
    runs = []
    for shape in slide.shapes:
        for r in shape._element.iter(f"{{{NS_A}}}r"):
            t = r.find(f"{{{NS_A}}}t")
            if t is not None and (t.text or "").strip():
                runs.append(r)
    return runs


def is_deleted_slide(slide, slide_w, slide_h):
    """
    True if slide should be removed:
    - Has a red diagonal deletion line
    - Every text run is strikethrough
    - Has a comment saying 'remove/delete this slide'
    """
    if _slide_has_delete_comment(slide):
        return True
    if _has_deletion_line(slide, slide_w, slide_h):
        return True
    runs = _all_text_runs_with_content(slide)
    if runs and all(_is_strikethrough(r) for r in runs):
        return True
    return False


def remove_slide(prs, idx):
    xml_slides = prs.slides._sldIdLst
    slide = prs.slides[idx]
    for rId, rel in list(prs.slides.part.rels.items()):
        try:
            if rel.target_part == slide.part:
                prs.slides.part.drop_rel(rId)
                break
        except Exception:
            pass
    xml_slides.remove(xml_slides[idx])


# ── Rule 2: Slide Numbers ──────────────────────────────────────────────────────

def _find_slide_num_shape(slide, slide_w):
    for shape in slide.shapes:
        txt = _shape_text(shape).replace("\x0b", "\n")
        if txt and SLIDE_NUM_PATTERN.match(txt):
            return shape
    # Fallback: yellow + top-right corner
    for shape in slide.shapes:
        if shape.shape_type in (1, 17) and _is_yellow(_get_fill_hex(shape._element)):
            try:
                if (shape.left or 0) > slide_w * 0.70 and (shape.top or 0) < 700000:
                    return shape
            except Exception:
                pass
    return None


def _set_slide_number_text(shape, label):
    try:
        tf = shape.text_frame
        new_text = f"Slide-{label}"
        # Skip the box entirely when the number is already correct — avoids adding
        # explicit black color to boxes that were previously inheriting theme color.
        existing_text = tf.text.strip()
        if existing_text == new_text:
            return
        for para in tf.paragraphs:
            runs = list(para.runs)
            if runs:
                runs[0].text = new_text
                runs[0].font.color.rgb = RGBColor(0, 0, 0)  # always black
                for r in runs[1:]:
                    r.text = ""
                for extra in tf.paragraphs[1:]:
                    for r in extra.runs:
                        r.text = ""
                return
        if tf.paragraphs:
            tf.paragraphs[0].text = new_text
    except Exception as e:
        print(f"    Warning (set slide num): {e}")


def _add_slide_number_box(slide, label, slide_w):
    """Create a yellow Slide-N box at top-right if none exists."""
    left   = Emu(10363200)
    top    = Emu(0)
    width  = Emu(max(slide_w - 10363200, 457200))
    height = Emu(457200)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = False

    sp_pr = txBox._element.find(_qn("p:spPr"))
    if sp_pr is None:
        sp_pr = etree.SubElement(txBox._element, _qn("p:spPr"))
    for old in sp_pr.findall(_qn("a:solidFill")):
        sp_pr.remove(old)
    sf = etree.SubElement(sp_pr, f"{{{NS_A}}}solidFill")
    srgb = etree.SubElement(sf, f"{{{NS_A}}}srgbClr")
    srgb.set("val", "FFFF00")

    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = f"Slide-{label}"
    run.font.bold = True
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(0, 0, 0)
    print(f"  [R2]  Added 'Slide-{label}' yellow box on slide")


# ── Rule 3: Remove Strikethrough Runs ─────────────────────────────────────────

def remove_strikethrough(element, label):
    count = 0
    for para_el in element.iter(f"{{{NS_A}}}p"):
        to_rm = [r for r in para_el.findall(f"{{{NS_A}}}r") if _is_strikethrough(r)]
        for r in to_rm:
            t = r.find(f"{{{NS_A}}}t")
            txt = (t.text or "").strip() if t is not None else ""
            para_el.remove(r)
            if txt:
                print(f"  [R3]  {label}: removed strikethrough '{txt[:50]}'")
            count += 1
    return count


# ── Rule 4: Remove Highlights + Fix Font Color ────────────────────────────────

def _layout_ph_default_text_color(slide, ph_idx, ph_type):
    """
    Look up the slide's layout for a placeholder matching ph_idx (or ph_type),
    return (is_light_scheme, hex_or_None) describing its default run text color.
    is_light_scheme=True means the default color is a white/light scheme color (bg1/lt1/lt2).
    Returns (False, None) if we can't determine.
    """
    LIGHT_SCHEME = {"bg1", "lt1", "lt2"}
    NS_P_local = "http://schemas.openxmlformats.org/presentationml/2006/main"
    NS_A_local = "http://schemas.openxmlformats.org/drawingml/2006/main"

    try:
        layout = slide.slide_layout
    except Exception:
        return False, None

    for lshape in layout.shapes:
        sp = lshape._element
        nvSpPr = sp.find(f"{{{NS_P_local}}}nvSpPr")
        if nvSpPr is None:
            continue
        nvPr = nvSpPr.find(f"{{{NS_P_local}}}nvPr")
        if nvPr is None:
            continue
        ph = nvPr.find(f"{{{NS_P_local}}}ph")
        if ph is None:
            continue
        if ph_idx is not None and ph.get("idx") != ph_idx:
            continue
        if ph_idx is None and ph_type is not None and ph.get("type") != ph_type:
            continue
        # Found matching layout placeholder
        txBody = sp.find(f"{{{NS_P_local}}}txBody")
        if txBody is None:
            break
        lstStyle = txBody.find(f"{{{NS_A_local}}}lstStyle")
        if lstStyle is None:
            break
        for lvl in lstStyle:
            for defRPr in lvl.findall(f"{{{NS_A_local}}}defRPr"):
                for sf in defRPr.iter(f"{{{NS_A_local}}}solidFill"):
                    for sc in sf.iter(f"{{{NS_A_local}}}schemeClr"):
                        val = sc.get("val", "")
                        return val in LIGHT_SCHEME, None
                    for srgb in sf.iter(f"{{{NS_A_local}}}srgbClr"):
                        v = srgb.get("val", "")
                        if len(v) == 6:
                            return False, v.upper()
        break
    return False, None


def _slide_has_dark_bg(slide, slide_w=None, slide_h=None):
    """
    Returns True if the slide (or its layout/master) has a dark background.
    Checks two mechanisms AMD templates use:
      1. <p:bg><a:bgPr> — explicit background fill (dark solidFill or blipFill)
      2. A full-slide <p:pic> with blipFill on the layout/master (AMD dark templates
         use a picture shape covering the entire slide as the background image)
    slide_w/slide_h must be provided (in EMU) for the full-slide pic check.
    """
    NS_P_local = "http://schemas.openxmlformats.org/presentationml/2006/main"
    NS_A_local = "http://schemas.openxmlformats.org/drawingml/2006/main"

    def _check_bg_element(el):
        """Check explicit <p:bg><*:bgPr> declarations.
        bgPr can be in either NS_P (layout level) or NS_A (slide level) depending on context.
        """
        for bg in el.iter(f"{{{NS_P_local}}}bg"):
            for bgPr in (bg.findall(f"{{{NS_P_local}}}bgPr") +
                         bg.findall(f"{{{NS_A_local}}}bgPr")):
                h = _get_fill_hex(bgPr)
                if h and _lum(h) < 128:
                    return True
                # blipFill means an image background — AMD templates always use dark images
                if (bgPr.find(f"{{{NS_A_local}}}blipFill") is not None or
                        bgPr.find(f"{{{NS_P_local}}}blipFill") is not None):
                    return True
        return False

    def _check_fullslide_pic(el):
        """
        Check for a <p:pic> with blipFill that covers most of the slide.
        AMD dark templates place a full-slide background image as a <p:pic>
        in the layout's spTree rather than in <p:bg>.
        Inside <p:pic>, both blipFill and spPr are under the NS_P namespace.
        """
        if slide_w is None or slide_h is None:
            return False
        for pic in el.iter(f"{{{NS_P_local}}}pic"):
            # blipFill is under p: namespace inside p:pic
            blip = pic.find(f"{{{NS_P_local}}}blipFill")
            if blip is None:
                continue
            # spPr is also under p: namespace inside p:pic
            sp_pr = pic.find(f"{{{NS_P_local}}}spPr")
            if sp_pr is None:
                continue
            # xfrm/ext inside spPr are under a: (drawingml)
            xfrm = sp_pr.find(f"{{{NS_A_local}}}xfrm")
            if xfrm is None:
                continue
            ext = xfrm.find(f"{{{NS_A_local}}}ext")
            if ext is None:
                continue
            try:
                cx = abs(int(ext.get("cx", 0)))
                cy = abs(int(ext.get("cy", 0)))
                # Covers ≥70% of slide in both dimensions → background image
                if cx >= slide_w * 0.70 and cy >= slide_h * 0.70:
                    return True
            except Exception:
                pass
        return False

    try:
        for obj in [slide, slide.slide_layout, slide.slide_layout.slide_master]:
            if _check_bg_element(obj._element):
                return True
            if _check_fullslide_pic(obj._element):
                return True
    except Exception:
        pass
    return False


def _get_shape_fill_hex(shape_el):
    """
    Return the hex fill color of the shape itself (not a child run).
    Looks inside <p:sp><p:spPr> solidFill → srgbClr.
    Returns None if no explicit solid fill.
    """
    NS_P_local = "http://schemas.openxmlformats.org/presentationml/2006/main"
    # <p:spPr> or <p:grpSpPr>
    for sp_pr in shape_el:
        tag = sp_pr.tag.split("}")[-1] if "}" in sp_pr.tag else sp_pr.tag
        if tag in ("spPr", "grpSpPr"):
            return _get_fill_hex(sp_pr)
    return None


def _get_run_color_hex_full(run_el):
    """
    Like _get_run_color_hex but also resolves schemeClr via AMD_SCHEME_COLORS.
    Returns hex string or None.
    """
    rPr = run_el.find(f"{{{NS_A}}}rPr")
    if rPr is None:
        return None
    for sf in rPr.iter(f"{{{NS_A}}}solidFill"):
        for srgb in sf.iter(f"{{{NS_A}}}srgbClr"):
            v = srgb.get("val", "")
            if len(v) == 6:
                return v.upper()
        for sc in sf.iter(f"{{{NS_A}}}schemeClr"):
            val = sc.get("val", "")
            if val in AMD_SCHEME_COLORS:
                return AMD_SCHEME_COLORS[val]
    return None


def _ambient_run_color(scope_el):
    """
    Return the most-common explicit font color among non-highlighted,
    non-strikethrough runs within scope_el (a paragraph or shape element).
    Resolves both srgbClr and schemeClr (via AMD_SCHEME_COLORS).
    Returns a hex string, or None if no explicit colors found.
    Red colors (#FF0000-family) are excluded — they are reviewer markup.
    """
    color_counts = {}
    for run_el in scope_el.iter(f"{{{NS_A}}}r"):
        rPr = run_el.find(f"{{{NS_A}}}rPr")
        if rPr is not None:
            if rPr.find(f"{{{NS_A}}}highlight") is not None:
                continue  # skip the highlighted runs we're trying to fix
            if rPr.get("strike", "noStrike") in ("sngStrike", "dblStrike"):
                continue  # skip struck-out markup runs
        color = _get_run_color_hex_full(run_el)
        if color and not _is_red(color):
            color_counts[color] = color_counts.get(color, 0) + 1
    if not color_counts:
        return None
    return max(color_counts, key=color_counts.get)


def _set_run_color(rPr, hex6):
    """
    Set the explicit font color on an rPr element.

    OOXML schema requires solidFill to appear before highlight, latin, ea, cs,
    sym, hlinkClick, etc. inside <a:rPr>. Using SubElement would append it at
    the end (after latin/highlight), violating the schema and triggering
    PowerPoint's repair dialog. We insert at position 0 (or after <a:ln> if
    present) to stay schema-compliant.
    """
    sf = rPr.find(f"{{{NS_A}}}solidFill")
    if sf is None:
        sf = etree.Element(f"{{{NS_A}}}solidFill")
        # solidFill must come after <a:ln> (if present) but before everything else
        ln = rPr.find(f"{{{NS_A}}}ln")
        children = list(rPr)
        insert_pos = children.index(ln) + 1 if ln is not None else 0
        rPr.insert(insert_pos, sf)
    for child in list(sf):
        sf.remove(child)
    srgb_el = etree.SubElement(sf, f"{{{NS_A}}}srgbClr")
    srgb_el.set("val", hex6)


def remove_highlights_fix_color(element, label, slide=None, slide_w=None, slide_h=None):
    """
    For every run that has a <a:highlight>:
    1. Remove the <a:highlight> element.
    2. Set font color to match the color of surrounding non-highlighted text,
       using this priority:
       P1. Sibling runs in the same paragraph have an explicit color → use it.
       P2. Other runs in the same shape have an explicit color → use it.
       P3. Run's own existing explicit color (non-red) → keep it.
       P4. Shape has an explicit solid fill → white (dark) or black (light).
       P5. No shape fill → check slide/layout/master background → white or black.
    Red colors (#FF0000-family) are always skipped — they are reviewer markup.
    Only highlighted runs are touched; all other runs are left completely alone.
    Applies to both slide content and speaker notes.
    """
    count = 0

    for run_el in element.iter(f"{{{NS_A}}}r"):
        rPr = run_el.find(f"{{{NS_A}}}rPr")
        if rPr is None:
            continue

        hl = rPr.find(f"{{{NS_A}}}highlight")
        if hl is None:
            continue

        # Read existing color BEFORE removing the highlight (resolve schemeClr too)
        existing_color = _get_run_color_hex_full(run_el)

        rPr.remove(hl)
        count += 1

        t = run_el.find(f"{{{NS_A}}}t")
        txt = (t.text or "") if t is not None else ""

        # Walk up to find the parent paragraph and parent shape
        para_ancestor = None
        sp_ancestor   = None
        anc = run_el.getparent()
        while anc is not None:
            tag = anc.tag.split("}")[-1] if "}" in anc.tag else anc.tag
            if tag == "p" and para_ancestor is None:
                para_ancestor = anc
            if tag == "sp":
                sp_ancestor = anc
                break
            anc = anc.getparent()

        # P1: color of non-highlighted sibling runs in the same paragraph
        if para_ancestor is not None:
            ambient = _ambient_run_color(para_ancestor)
            if ambient:
                _set_run_color(rPr, ambient)
                print(f"  [R4]  {label}: highlight removed + color matched #{ambient} (para siblings) for '{txt[:40]}'")
                continue

        # P2: color of non-highlighted runs anywhere in the same shape
        if sp_ancestor is not None:
            ambient = _ambient_run_color(sp_ancestor)
            if ambient:
                _set_run_color(rPr, ambient)
                print(f"  [R4]  {label}: highlight removed + color matched #{ambient} (shape siblings) for '{txt[:40]}'")
                continue

        # P3: run's own existing explicit color (skip red markup colors)
        if existing_color and not _is_red(existing_color):
            _set_run_color(rPr, existing_color)
            print(f"  [R4]  {label}: highlight removed + color preserved #{existing_color} for '{txt[:40]}'")
            continue

        # P4: infer from shape solid fill
        shape_fill_hex = _get_shape_fill_hex(sp_ancestor) if sp_ancestor is not None else None
        if shape_fill_hex:
            target = "FFFFFF" if _lum(shape_fill_hex) < 128 else "000000"
            color_name = "white" if target == "FFFFFF" else "black"
            _set_run_color(rPr, target)
            print(f"  [R4]  {label}: highlight removed + color→{color_name} (shape fill #{shape_fill_hex}) for '{txt[:40]}'")
            continue

        # P5: infer from slide/layout/master background
        if slide is not None and _slide_has_dark_bg(slide, slide_w=slide_w, slide_h=slide_h):
            _set_run_color(rPr, "FFFFFF")
            print(f"  [R4]  {label}: highlight removed + color→white (dark slide bg) for '{txt[:40]}'")
        else:
            _set_run_color(rPr, "000000")
            print(f"  [R4]  {label}: highlight removed + color→black for '{txt[:40]}'")

    return count


# ── Rule 4b: Fix font color for all text in filled shapes ─────────────────────

def fix_filled_shape_font_color(slide, sn):
    """
    For every shape with an explicit solid fill (including AMD theme colors),
    set all text runs to white (dark fill, lum < 128) or black (light fill).
    This covers text that has no highlight tag and was missed by Rule 4.
    Slide-number yellow boxes are unaffected (their text is already black from Rule 2).
    """
    count = 0
    for shape in slide.shapes:
        if not (hasattr(shape, "has_text_frame") and shape.has_text_frame):
            continue
        fill_hex = _get_shape_fill_hex(shape._element)
        if not fill_hex:
            continue
        lum = _lum(fill_hex)
        target = "FFFFFF" if lum < 128 else "000000"
        for run_el in shape._element.iter(f"{{{NS_A}}}r"):
            t = run_el.find(f"{{{NS_A}}}t")
            txt = (t.text or "")[:25] if t is not None else ""
            if not txt.strip():
                continue  # skip empty runs
            rPr = run_el.find(f"{{{NS_A}}}rPr")
            if rPr is None:
                # Create rPr — must be first child of <a:r>
                rPr = etree.Element(f"{{{NS_A}}}rPr")
                run_el.insert(0, rPr)
            current = _get_run_color_hex(run_el)
            if current and current.upper() == target:
                continue  # already correct
            _set_run_color(rPr, target)
            color_name = "white" if lum < 128 else "black"
            print(f"  [R4b] Slide {sn}: fill #{fill_hex} → {color_name} for '{txt}'")
            count += 1
    return count


# ── Rule 5: Clean empty paragraphs after strikeout removal ────────────────────

def clean_empty_paragraphs(slide_obj, max_consecutive=0):
    """
    Remove excess empty paragraphs from all text frames in slide_obj.
    max_consecutive=0 → remove ALL empty paragraphs (used for slide content).
    max_consecutive=1 → collapse runs of blanks to one (used for speaker notes).
    The mandatory last paragraph in each text frame is never removed.
    """
    count = 0
    for shape in slide_obj.shapes:
        if not (hasattr(shape, "has_text_frame") and shape.has_text_frame):
            continue
        tf = shape.text_frame
        paras = list(tf.paragraphs)
        if len(paras) <= 1:
            continue
        to_rm = []
        consecutive_empty = 0
        for para in paras[:-1]:  # never touch the last paragraph
            if not para.text.strip():
                consecutive_empty += 1
                if consecutive_empty > max_consecutive:
                    to_rm.append(para._p)
                    count += 1
            else:
                consecutive_empty = 0
        for p_el in to_rm:
            parent = p_el.getparent()
            if parent is not None:
                parent.remove(p_el)
    return count


# ── Rules 6+9: Remove ALL shapes outside the slide (except NTD/Shared yellow) ─

def _is_outside_slide(shape, slide_w):
    """True if shape is placed fully or mostly outside the left/right slide boundary."""
    try:
        left  = shape.left  or 0
        width = shape.width or 0
        right = left + width
        if right <= 0:
            return True
        if left >= slide_w:
            return True
        if left < 0 and (-left) >= width * 0.80:
            return True
        if right > slide_w and (right - slide_w) >= width * 0.80:
            return True
    except Exception:
        pass
    return False


def _is_ntd_protected_box(shape):
    """
    True if this yellow NTD/Shared box should be kept outside the slide.
    Protection is lifted when the box has BOTH strikethrough AND red-colored
    text runs — that signals deleted/deprecated content and the box is removed.
    """
    if not _is_yellow(_get_fill_hex(shape._element)):
        return False
    txt = _shape_text(shape)
    if not NTD_PATTERN.search(txt):
        return False
    # Lift protection if the box contains any struck-out+red runs
    all_runs = [r for r in shape._element.iter(f"{{{NS_A}}}r")
                if (r.find(f"{{{NS_A}}}t") is not None
                    and (r.find(f"{{{NS_A}}}t").text or "").strip())]
    if all_runs:
        has_strike = any(_is_strikethrough(r) for r in all_runs)
        has_red    = any(_is_red(_get_run_color_hex(r) or "") for r in all_runs)
        if has_strike and has_red:
            return False  # struck+red content → remove despite NTD label
    return True


def _should_remove_outside_shape(shape, slide_w):
    """
    Rules 6+9 combined: remove any shape outside the slide boundary,
    EXCEPT yellow NTD/Shared boxes which are always kept.
    """
    if not _is_outside_slide(shape, slide_w):
        return False
    if _is_ntd_protected_box(shape):
        return False
    return True


# ── Rule 8: Remove double spaces + leading spaces per paragraph ───────────────

def remove_double_spaces(element):
    count = 0
    for t_el in element.iter(f"{{{NS_A}}}t"):
        if not (t_el.text and "  " in t_el.text):
            continue
        # Do not collapse spaces inside NTD / Fully Shared / Partially Shared boxes
        sp_ancestor = t_el
        while sp_ancestor is not None:
            tag = sp_ancestor.tag.split("}")[-1] if "}" in sp_ancestor.tag else sp_ancestor.tag
            if tag == "sp":
                break
            sp_ancestor = sp_ancestor.getparent()
        if sp_ancestor is not None:
            sp_txt = "".join(
                (t.text or "")
                for t in sp_ancestor.iter(f"{{{NS_A}}}t")
            )
            if NTD_PATTERN.search(sp_txt):
                continue
        t_el.text = re.sub(r" {2,}", " ", t_el.text)
        count += 1
    return count


def trim_paragraph_leading_spaces(element):
    """
    Trim leading spaces from the first text run of each paragraph.
    After strikethrough removal, a deleted run at the start of a paragraph
    can leave the next run beginning with a space, creating a visual indent.
    """
    count = 0
    for para_el in element.iter(f"{{{NS_A}}}p"):
        for run_el in para_el.findall(f"{{{NS_A}}}r"):
            t = run_el.find(f"{{{NS_A}}}t")
            if t is not None and t.text:
                stripped = t.text.lstrip(" ")
                if stripped != t.text:
                    t.text = stripped
                    count += 1
                break  # only the first run per paragraph
    return count


# ── Rule 9: Remove ALL blue boxes outside slide ───────────────────────────────

def _shape_has_blue(shape_el):
    """True if the shape has a blue solidFill (rgb or theme color)."""
    # Check solid fill via srgbClr
    fill = _get_fill_hex(shape_el)
    if fill and _is_blue(fill):
        return True
    # Check solid fill via theme schemeClr (accent1, accent2 etc. resolve to blue
    # in the AMD/standard Office theme)
    for sf in shape_el.iter(f"{{{NS_A}}}solidFill"):
        for sc in sf.iter(f"{{{NS_A}}}schemeClr"):
            val = sc.get("val", "")
            # accent1 = AMD red/blue depending on theme; accent2 = blue in Office default
            # We conservatively flag both accent1 and accent2 when outside the slide
            if val in ("accent1", "accent2", "dk1", "hyperlink"):
                return True
    return False


def _is_blue_box_outside(shape, slide_w):
    if not _shape_has_blue(shape._element):
        return False
    return _is_outside_slide(shape, slide_w)


# ── Rule 14: Remove empty yellow boxes ────────────────────────────────────────

def _is_empty_yellow_box(shape):
    """
    Rule 14: yellow boxes whose visible text is empty (content was deleted by
    strikethrough removal or was never added). Slide-number boxes are never
    empty at this point (text was set in Rule 2), so they are not affected.
    """
    if not _is_yellow(_get_fill_hex(shape._element)):
        return False
    return not _shape_text(shape).strip()


# ── Rule 13: Animation notes yellow box removal ───────────────────────────────

def _is_animation_notes_box(shape):
    """
    Rule 13: Remove yellow-filled boxes whose text begins with 'Animation'
    (authoring cues written for developers / animators, not live content).
    NTD/Shared boxes are NOT animation notes and are never matched.
    """
    if not _is_yellow(_get_fill_hex(shape._element)):
        return False
    txt = _shape_text(shape)
    if not txt:
        return False
    return bool(ANIMATION_NOTES_PATTERN.match(txt))


# ── Rules 10 & 11: Yellow box removal ─────────────────────────────────────────

def _is_yellow_strikeout_red_box(shape):
    """
    Rule 10: Remove yellow box if it has BOTH strikethrough AND red-colored text.
    Rule 11: Protect NTD / Shared boxes ONLY when their text is NOT struck-out+red.
             If the NTD/Shared text is itself strikethrough+red, it is marked for
             deletion and should be removed.
    Logic:
      - If every non-empty run is strikethrough+red → remove regardless of NTD text
        (the whole box content is struck-out authoring markup)
      - If box has NTD/Shared text that has live (non-struck) runs → protect it
      - Otherwise: remove if has_strike AND has_red
    """
    if not _is_yellow(_get_fill_hex(shape._element)):
        return False
    txt = _shape_text(shape)
    if not txt:
        return False

    all_runs = [r for r in shape._element.iter(f"{{{NS_A}}}r")
                if (r.find(f"{{{NS_A}}}t") is not None
                    and (r.find(f"{{{NS_A}}}t").text or "").strip())]
    if not all_runs:
        return False

    # Classify each run
    strike_runs = [r for r in all_runs if _is_strikethrough(r)]
    red_runs    = [r for r in all_runs
                   if _is_red(_get_run_color_hex(r) or "")]
    live_runs   = [r for r in all_runs if not _is_strikethrough(r)]

    has_strike = bool(strike_runs)
    has_red    = bool(red_runs)

    if not (has_strike and has_red):
        return False  # doesn't meet Rule 10 criteria at all

    # If there are live (non-strikethrough) runs → check Rule 11 protection
    if live_runs and NTD_PATTERN.search(txt):
        return False  # active NTD/Shared informational box — keep it

    # All content is struck-out/red, or no live NTD text → remove
    return True


# ── Main ───────────────────────────────────────────────────────────────────────

def make_output_path(input_path):
    """Rule 12: strip _V<n> version token from filename."""
    import os
    dirname  = os.path.dirname(input_path)
    basename = os.path.basename(input_path)
    new_name = VERSION_PATTERN.sub("", basename)
    if new_name == basename:
        # No version token found — add _Cleaned to avoid overwriting
        stem, ext = os.path.splitext(basename)
        new_name = stem + "_Cleaned" + ext
    return os.path.join(dirname, new_name)


def cleanup(input_path, output_path):
    prs     = Presentation(input_path)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    print(f"Input : {input_path}")
    print(f"Output: {output_path}")
    print(f"Slides: {len(prs.slides)}  |  {slide_w} x {slide_h} EMU\n")

    # ── Rule 7: Remove deleted slides first (reverse order) ───────────────────
    slides_removed  = 0
    deleted_indices = []
    for idx in range(len(prs.slides)):
        s = prs.slides[idx]
        if _slide_has_delete_comment(s):
            deleted_indices.append(idx)
            print(f"  [R7]  Slide {idx+1}: marked for removal (comment says delete/remove slide)")
        elif is_deleted_slide(s, slide_w, slide_h):
            deleted_indices.append(idx)
            print(f"  [R7]  Slide {idx+1}: marked for removal (diagonal deletion line or all-strikethrough)")

    for idx in reversed(deleted_indices):
        remove_slide(prs, idx)
        slides_removed += 1
    print(f"\n  {slides_removed} slide(s) removed. Remaining: {len(prs.slides)}\n")

    slides_list = list(prs.slides)

    # ── Identify Learning Objectives slide for Rule 2 ─────────────────────────
    lo_idx = None
    for idx, slide in enumerate(slides_list):
        if LO_PATTERN.search(_slide_title(slide)):
            lo_idx = idx
            print(f"  [R2]  LO slide at index {idx+1} → counted as Slide-2 (no label shown); "
                  f"slide after → Slide-3\n")
            break

    # ── Per-slide processing ───────────────────────────────────────────────────
    total_shapes = 0
    current_num  = 1       # will be set to 2 at LO (no label), 3 at slide-after-LO
    current_title = None

    for idx, slide in enumerate(slides_list):
        sn = idx + 1  # 1-based for logging

        # Rule 1
        remove_comments(slide, sn)

        # Rule 2: compute label
        title = _slide_title(slide)
        slide_label = None

        if lo_idx is not None:
            if idx < lo_idx:
                slide_label = None          # title slide — no yellow number
            elif idx == lo_idx:
                current_num   = 2           # LO counted as Slide-2 but no label shown
                slide_label   = None
                current_title = title
                # Remove any existing slide-number box on the LO slide
                sn_shape = _find_slide_num_shape(slide, slide_w)
                if sn_shape:
                    sn_shape._element.getparent().remove(sn_shape._element)
                    print(f"  [R2]  Slide {sn} (LO): removed existing slide-number box")
            elif idx == lo_idx + 1:
                # Slide immediately after LO starts at Slide-3
                current_num   = 3
                slide_label   = 3
                current_title = title
            else:
                same = bool(title and title == current_title)
                # "Apply Your Knowledge" slides always get their own number
                # even when back-to-back with the same title
                is_ayk = bool(AYK_PATTERN.search(title or ""))
                if same and not is_ayk:
                    slide_label = current_num
                else:
                    current_num  += 1
                    slide_label   = current_num
                    current_title = title

        if slide_label is not None:
            sn_shape = _find_slide_num_shape(slide, slide_w)
            if sn_shape:
                old = _shape_text(sn_shape)
                norm_old = old.strip().replace(" ", "-").upper()
                norm_new = f"SLIDE-{slide_label}"
                _set_slide_number_text(sn_shape, slide_label)
                if norm_old != norm_new:
                    print(f"  [R2]  Slide {sn}: '{old}' → 'Slide-{slide_label}'")

        # Collect shapes to remove (Rules 6, 9, 10+11)
        to_remove = []

        for shape in slide.shapes:
            # Rule 13: animation notes yellow box
            if _is_animation_notes_box(shape):
                snippet = _shape_text(shape)[:55].replace("\n", " ")
                print(f"  [R13] Slide {sn}: removing animation notes box '{shape.name}' → '{snippet}'")
                to_remove.append(shape)
                continue

            # Rule 10+11: yellow box with strikeout + red, not NTD/Shared
            if _is_yellow_strikeout_red_box(shape):
                snippet = _shape_text(shape)[:55].replace("\n", " ")
                print(f"  [R10] Slide {sn}: removing yellow box '{shape.name}' → '{snippet}'")
                to_remove.append(shape)
                continue

            # Rules 6+9: any shape outside the slide (except NTD/Shared yellow boxes)
            if _should_remove_outside_shape(shape, slide_w):
                print(f"  [R6]  Slide {sn}: removing '{shape.name}' (outside slide)")
                to_remove.append(shape)
                continue

        for shape in to_remove:
            try:
                shape._element.getparent().remove(shape._element)
                total_shapes += 1
            except Exception as e:
                print(f"    Warning: cannot remove '{shape.name}': {e}")

        # Rule 3: strikethrough runs
        remove_strikethrough(slide._element, f"Slide {sn}")

        # Rule 4: remove highlights + fix light-colored fonts
        remove_highlights_fix_color(slide._element, f"Slide {sn}", slide=slide, slide_w=slide_w, slide_h=slide_h)

        # Rule 5: clean empty paragraphs
        clean_empty_paragraphs(slide)

        # Rule 8: double spaces + leading spaces per paragraph
        remove_double_spaces(slide._element)
        trim_paragraph_leading_spaces(slide._element)

        # Rule 14: second pass — remove yellow boxes emptied by Rule 3
        for shape in list(slide.shapes):
            if _is_empty_yellow_box(shape):
                print(f"  [R14] Slide {sn}: removing empty yellow box '{shape.name}'")
                try:
                    shape._element.getparent().remove(shape._element)
                    total_shapes += 1
                except Exception as e:
                    print(f"    Warning: cannot remove '{shape.name}': {e}")

        # Notes: Rules 3, 4, 5, 8
        try:
            if slide.has_notes_slide:
                ns_el = slide.notes_slide._element
                remove_strikethrough(ns_el, f"Slide {sn} Notes")
                remove_highlights_fix_color(ns_el, f"Slide {sn} Notes", slide=slide, slide_w=slide_w, slide_h=slide_h)
                clean_empty_paragraphs(slide.notes_slide, max_consecutive=1)
                remove_double_spaces(ns_el)
                trim_paragraph_leading_spaces(ns_el)
        except Exception:
            pass

    print(f"\n{'─'*62}")
    print(f"Shapes removed : {total_shapes}")
    print(f"Slides removed : {slides_removed}")
    print(f"Saving → {output_path}")
    prs.save(output_path)
    print("Done.")


if __name__ == "__main__":
    if len(sys.argv) < 2 or len(sys.argv) > 3:
        print("Usage: python cleanup_sb.py <input.pptx> [output.pptx]")
        sys.exit(1)
    inp = sys.argv[1]
    out = sys.argv[2] if len(sys.argv) == 3 else make_output_path(inp)
    cleanup(inp, out)
