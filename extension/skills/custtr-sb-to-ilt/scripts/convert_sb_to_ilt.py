"""
SB to ILT Converter — bundled skill script
Usage: python convert_sb_to_ilt.py <input.pptx> <output.pptx>

Applies the ILT conversion rules:
  1. Remove slide-number / branch-label annotations (yellow patches, top-right corner)
  2. Remove boxes placed outside the standard slide layout — by NTD/authoring text OR
     by position+color (FFFF00 fill outside slide bounds or at top-right corner)
  3. Remove cross (X) close button from popup/overlay slides (small picture at top-right,
     NOT on Apply Your Knowledge / Knowledge Check slides)
  4. Remove click/drag/interaction instruction boxes (entire shape or entire parent group,
     including the accompanying (i) info icon) from slides and voiceover notes.
     Preserves "Select all that apply".
  5. Remove orange curved arrows, including SVG-based ones ("Back outline")
"""

import re
import sys
from pptx import Presentation
from lxml import etree

# ── Rule 1: Slide-number / branch-label patterns ──────────────────────────────
# Matches: "Slide-3", "Slide 05", "Slide 05 Branch 1.1", "Slide 05Branch 1.1",
#          "New Slide 06", "Slide 06" and similar authoring annotations.
SLIDE_NUM_PATTERN = re.compile(
    r"^(Slide[-\s]\d+(\s*([\n\r\x0b/\\]\s*)?Branch\s*[\d.]+)?|New\s*[\n\r\x0b]?\s*Slide\s+\d+|Slide\s+\d+)$",
    re.IGNORECASE | re.DOTALL,
)

# ── Rule 2: Outside-layout box patterns ───────────────────────────────────────
# Primary: text-based (NTD annotations, shared-slide markers)
OUTSIDE_LAYOUT_PATTERN = re.compile(
    r"(NTD\s*:|Fully Shared Slide|Partially Shared Slide|Shared Slide)",
    re.IGNORECASE,
)
# Secondary: yellow fill (FFFF00) positioned outside or at top-right — catches
# authoring boxes that have no recognisable text marker.
YELLOW_HEX = "FFFF00"

# ── Rule 4: Click/interaction instruction patterns ────────────────────────────
# NOTE: "Select all that apply" is intentionally NOT in this list — keep it.
CLICK_PATTERNS = [
    re.compile(r"click\s+each", re.IGNORECASE),
    re.compile(r"click\s+example", re.IGNORECASE),
    re.compile(r"click\s+through", re.IGNORECASE),
    re.compile(r"click\s+the\s+arrow", re.IGNORECASE),
    re.compile(r"\bclick\s+\w+", re.IGNORECASE),
    # "Drag-and-drop" (literal) AND "Drag X and drop it" / "Drag a step ... drop"
    re.compile(r"drag[- ]and[- ]drop", re.IGNORECASE),
    re.compile(r"\bdrag\b.{0,60}\bdrop\b", re.IGNORECASE | re.DOTALL),
    re.compile(r"match\s+(each|the)", re.IGNORECASE),
    re.compile(r"hover\s+over", re.IGNORECASE),
    re.compile(r"roll\s+over", re.IGNORECASE),
    re.compile(r"tap\s+on", re.IGNORECASE),
]

# Sentence-level removal for voiceover notes
CLICK_SENTENCE_PATTERN = re.compile(
    r"[^.!?]*(?:click\s+each|click\s+example|click\s+through|click\s+the\s+arrow"
    r"|drag[- ]and[- ]drop|\bdrag\b.{0,60}\bdrop\b|match\s+(?:each|the)"
    r"|hover\s+over|roll\s+over|tap\s+on)[^.!?]*[.!?]?",
    re.IGNORECASE,
)

# Short filler phrases that may accompany click instructions
FILLER_PATTERN = re.compile(
    r"^(\.{2,}\s*)?(to\s+)?(learn\s+more|view\s+more|see\s+more|find\s+out\s+more)\.?$",
    re.IGNORECASE,
)

# ── Quiz / popup slide keywords ────────────────────────────────────────────────
QUIZ_KW = {"apply your knowledge", "knowledge check"}

# PPTX XML namespaces
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


# ── Helpers ───────────────────────────────────────────────────────────────────

def get_shape_text(shape):
    try:
        if shape.has_text_frame:
            return "\n".join(p.text for p in shape.text_frame.paragraphs).strip()
    except Exception:
        pass
    return ""


def get_element_text(el):
    """Get all text from any XML element recursively."""
    return "".join((t.text or "") for t in el.iter(f"{{{NS_A}}}t")).strip()


def normalise(txt):
    """Replace PPTX paragraph-break character with newline."""
    return txt.replace("\x0b", "\n")


def is_click_text(txt):
    """Returns True if the text matches any click/interaction pattern."""
    return any(p.search(txt) for p in CLICK_PATTERNS)


def is_pure_click_text(txt):
    """
    Returns True if ALL non-empty lines of txt are click instructions or filler.
    Used to decide whether to remove a shape entirely.
    """
    if not txt.strip():
        return False
    if not is_click_text(txt):
        return False
    lines = [l.strip() for l in normalise(txt).splitlines() if l.strip()]
    return all(
        is_click_text(line) or FILLER_PATTERN.match(line)
        for line in lines
    )


# ── Yellow fill helper (Rules 1 & 2 fallback) ────────────────────────────────

def has_yellow_fill(shape):
    """Returns True if any srgbClr in the shape element is FFFF00."""
    NS_A_FULL = f"{{{NS_A}}}"
    for clr in shape._element.iter(f"{NS_A_FULL}srgbClr"):
        if clr.get("val", "").upper() == YELLOW_HEX:
            return True
    return False


# ── Rule 1 ────────────────────────────────────────────────────────────────────

def is_slide_number_shape(shape, slide_w=12192000):
    """
    Removes yellow annotation boxes at the top-right corner of slides.
    Primary: text matches slide-number / branch-label pattern.
    Fallback: FFFF00 fill + positioned in the top-right quadrant AND within
    the slide boundary (catches non-standard labels like "Slide 05 Branch 1.1",
    "New Slide 06", etc. without stealing from Rule 2's outside-layout detection).
    """
    txt = get_shape_text(shape)
    # Primary: text pattern match
    if txt and SLIDE_NUM_PATTERN.match(normalise(txt)):
        return True
    # Fallback: yellow fill at top-right, within slide bounds
    # Only for rectangle/text shapes (shape_type 1=rectangle, 17=text box)
    if shape.shape_type in (1, 17) and has_yellow_fill(shape):
        try:
            left = shape.left or 0
            top = shape.top or 0
            # Must be inside the slide (not an outside-layout box),
            # in the right 25% and top 10% — the annotation corner.
            if left <= slide_w and left > 9000000 and top < 700000:
                return True
        except Exception:
            pass
    return False


# ── Rule 2 ────────────────────────────────────────────────────────────────────

def is_outside_layout_box(shape, slide_w=12192000):
    """
    Removes authoring boxes that sit outside the visible slide area.
    Primary: text contains NTD/Shared Slide markers.
    Fallback: shape left-edge exceeds slide width AND has yellow fill —
    catches NTD boxes that lack standard text markers.
    """
    txt = get_shape_text(shape)
    if txt and OUTSIDE_LAYOUT_PATTERN.search(txt):
        return True
    # Fallback: positioned beyond slide right edge with yellow fill
    try:
        if (shape.left or 0) > slide_w and has_yellow_fill(shape):
            return True
    except Exception:
        pass
    return False


# ── Rule 3 ────────────────────────────────────────────────────────────────────

def is_quiz_slide(slide):
    """Detect Apply Your Knowledge / Knowledge Check slides."""
    for shape in slide.shapes:
        try:
            if shape.has_text_frame:
                if any(kw in get_shape_text(shape).lower() for kw in QUIZ_KW):
                    return True
        except Exception:
            pass
    return False


def is_popup_slide(slide, slide_w, slide_h):
    """
    Popup/overlay slides contain a near-full-size background rectangle
    (covering > 90% width and > 90% height) that acts as the overlay panel.
    """
    for shape in slide.shapes:
        if shape.shape_type != 1:
            continue
        try:
            l = (shape.left or 0) / slide_w * 100
            t = (shape.top or 0) / slide_h * 100
            w = (shape.width or 0) / slide_w * 100
            h = (shape.height or 0) / slide_h * 100
            if l < 5 and t < 5 and w > 90 and h > 90:
                return True
        except Exception:
            pass
    return False


def is_popup_close_button(shape, slide, slide_w, slide_h):
    """
    Rule 3: Small Picture at top-right of a popup slide (non-quiz).
    Criteria: Picture type, left > 85%, top < 8%, width < 8%, height < 8%.
    """
    if shape.shape_type != 13:
        return False
    if is_quiz_slide(slide):
        return False
    if not is_popup_slide(slide, slide_w, slide_h):
        return False
    try:
        l = (shape.left or 0) / slide_w * 100
        t = (shape.top or 0) / slide_h * 100
        w = (shape.width or 0) / slide_w * 100
        h = (shape.height or 0) / slide_h * 100
        return l > 85 and t < 8 and w < 8 and h < 8
    except Exception:
        return False


# ── Rule 4 ────────────────────────────────────────────────────────────────────

def is_click_instruction_shape(shape):
    """
    For non-group top-level shapes: returns True if the shape's entire text
    is a click instruction so the whole shape should be removed.
    """
    if not shape.has_text_frame:
        return False
    return is_pure_click_text(normalise(get_shape_text(shape)))


def group_contains_click_instruction(shape):
    """
    For GROUP top-level shapes: returns True if any direct or nested <p:sp>
    within the group has pure click instruction text.
    When True, the ENTIRE group (including the (i) icon sub-group) is removed.
    """
    if shape.shape_type != 6:  # GROUP
        return False
    for sp in shape._element.iter(f"{{{NS_P}}}sp"):
        txt = get_element_text(sp).strip()
        if txt and is_pure_click_text(txt):
            return True
    return False


def slide_has_interactivity(slide):
    """
    Returns True if this slide is an interactivity slide — i.e. it contains
    a top-level group with a click/interaction instruction inside.
    Used by Rule 6 to limit voiceover cleaning to interactivity slides only.
    """
    for shape in slide.shapes:
        if shape.shape_type == 6:  # GROUP
            for sp in shape._element.iter(f"{{{NS_P}}}sp"):
                txt = get_element_text(sp).strip()
                if txt and is_click_text(txt):
                    return True
    return False


def clean_notes(notes_slide, slide_num):
    """
    Rule 6: Remove click-instruction sentences from voiceover notes,
    but ONLY on interactivity slides (caller is responsible for the guard).

    Strategy: check each paragraph's full text against CLICK_PATTERNS.
    - If the whole paragraph is a click instruction → remove the paragraph element.
    - If the paragraph has mixed content → strip only the click-instruction portion
      using CLICK_SENTENCE_PATTERN and leave the rest.
    """
    if notes_slide is None:
        return

    for ph in notes_slide.placeholders:
        if not ph.has_text_frame:
            continue
        tf = ph.text_frame
        paras_to_remove = []

        for para in tf.paragraphs:
            full_para = "".join(run.text for run in para.runs).strip()
            if not full_para:
                continue
            if not any(p.search(full_para) for p in CLICK_PATTERNS):
                continue

            # Check if the whole paragraph is purely a click instruction
            if is_pure_click_text(full_para):
                paras_to_remove.append(para._p)
                print(f"  [R6-VoiceOver]  Slide {slide_num}: removed paragraph: '{full_para[:70]}'")
            else:
                # Mixed content — strip only the matching sentence
                for run in para.runs:
                    original = run.text
                    cleaned = CLICK_SENTENCE_PATTERN.sub("", original)
                    cleaned = re.sub(r"\s{2,}", " ", cleaned).strip()
                    if cleaned != original:
                        run.text = cleaned
                        print(f"  [R6-VoiceOver]  Slide {slide_num}: cleaned sentence from notes")

        for p_el in paras_to_remove:
            parent = p_el.getparent()
            if parent is not None:
                parent.remove(p_el)


# ── Rule 5 ────────────────────────────────────────────────────────────────────

def is_svg_curved_arrow(shape):
    try:
        cNvPr = shape._element.find(f".//{{{NS_A}}}cNvPr")
        if cNvPr is not None and cNvPr.get("descr", "").strip().lower() == "back outline":
            return True
    except Exception:
        pass
    return False


def is_orange_curved_arrow(shape):
    NS = NS_A
    try:
        xml = shape._element
        prstGeom = xml.find(f".//{{{NS}}}prstGeom")
        if prstGeom is None:
            return False
        prst = prstGeom.get("prst", "").lower()
        if "curve" not in prst and "arc" not in prst:
            return False
        for srgbClr in xml.iter(f"{{{NS}}}srgbClr"):
            val = srgbClr.get("val", "")
            if len(val) != 6:
                continue
            r, g, b = int(val[0:2], 16), int(val[2:4], 16), int(val[4:6], 16)
            # Red: R>180, G<80, B<80  |  Orange: R>180, 80<=G<=200, B<80
            if r > 180 and g <= 200 and b < 80:
                return True
    except Exception:
        pass
    return False


def remove_svg_arrows_from_groups(slide, slide_num):
    """Rule 5: Remove SVG curved arrows nested inside group shapes."""
    removed = 0
    targets = []
    for el in slide._element.iter():
        tag = el.tag.split("}")[-1] if "}" in el.tag else el.tag
        if tag == "cNvPr" and el.get("descr", "").strip().lower() == "back outline":
            parent = el.getparent()
            while parent is not None:
                ptag = parent.tag.split("}")[-1] if "}" in parent.tag else parent.tag
                if ptag == "pic":
                    targets.append(parent)
                    break
                parent = parent.getparent()
    for pic_el in targets:
        grandparent = pic_el.getparent()
        if grandparent is not None:
            cn = pic_el.find(f".//{{{NS_A}}}cNvPr")
            name = cn.get("name", "?") if cn is not None else "?"
            grandparent.remove(pic_el)
            print(f"  [R5-SvgArrow]   Slide {slide_num}: '{name}' (nested SVG arrow removed)")
            removed += 1
    return removed


# ── Core removal helper ───────────────────────────────────────────────────────

def remove_shape(shape):
    shape._element.getparent().remove(shape._element)


# ── Main conversion ───────────────────────────────────────────────────────────

def convert(input_path, output_path):
    prs = Presentation(input_path)
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    total_removed = 0

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1

        # Rule 6: detect interactivity BEFORE shapes are removed this iteration
        is_interactive = slide_has_interactivity(slide)

        shapes_to_remove = []

        for shape in slide.shapes:
            # Rule 1: slide-number label in yellow patch
            if is_slide_number_shape(shape, slide_w):
                print(f"  [R1-SlideNum]   Slide {slide_num}: '{shape.name}' -> '{normalise(get_shape_text(shape))[:40]}'")
                shapes_to_remove.append(shape)
                continue

            # Rule 2: box outside the standard slide layout
            if is_outside_layout_box(shape, slide_w):
                snippet = get_shape_text(shape)[:50].replace("\n", " ")
                print(f"  [R2-OutsideBox] Slide {slide_num}: '{shape.name}' -> '{snippet}...'")
                shapes_to_remove.append(shape)
                continue

            # Rule 3: popup close button (X at top-right, non-quiz popup slide only)
            if is_popup_close_button(shape, slide, slide_w, slide_h):
                print(f"  [R3-CloseBtn]   Slide {slide_num}: '{shape.name}' (popup close button removed)")
                shapes_to_remove.append(shape)
                continue

            # Rule 4a: GROUP containing click instruction → remove entire group
            # (this also removes the (i) info icon sub-group inside it)
            if group_contains_click_instruction(shape):
                all_txt = normalise(get_element_text(shape._element))[:60].replace("\n", " ")
                print(f"  [R4-ClickGrp]   Slide {slide_num}: '{shape.name}' -> entire group removed (click+icon)")
                shapes_to_remove.append(shape)
                continue

            # Rule 4b: standalone click-instruction shape → remove entire shape
            if is_click_instruction_shape(shape):
                snippet = normalise(get_shape_text(shape))[:60].replace("\n", " ")
                print(f"  [R4-ClickBox]   Slide {slide_num}: '{shape.name}' -> '{snippet}'")
                shapes_to_remove.append(shape)
                continue

            # Rule 5a: geometric red or orange curved arrow
            if is_orange_curved_arrow(shape):
                print(f"  [R5-RedOrangeArrow] Slide {slide_num}: '{shape.name}'")
                shapes_to_remove.append(shape)
                continue

            # Rule 5b: top-level SVG curved arrow ("Back outline")
            if is_svg_curved_arrow(shape):
                print(f"  [R5-SvgArrow]   Slide {slide_num}: '{shape.name}'")
                shapes_to_remove.append(shape)
                continue

        for shape in shapes_to_remove:
            remove_shape(shape)
        total_removed += len(shapes_to_remove)

        # Rule 5: nested SVG arrows inside groups not yet removed
        total_removed += remove_svg_arrows_from_groups(slide, slide_num)

        # Rule 6: clean voiceover notes — interactivity slides only
        # (use is_interactive captured BEFORE shapes were removed this iteration)
        try:
            if slide.has_notes_slide and is_interactive:
                clean_notes(slide.notes_slide, slide_num)
        except Exception:
            pass

    print(f"\nDone. {total_removed} shapes removed.")
    prs.save(output_path)
    print(f"Saved: {output_path}")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python convert_sb_to_ilt.py <input.pptx> <output.pptx>")
        sys.exit(1)
    convert(sys.argv[1], sys.argv[2])
