"""
Extract slide content from a .pptx file for source-finding.
Outputs JSON with per-slide text and image descriptions.

Usage:
    python extract_slides.py <path_to_pptx>
"""

import sys
import re
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Authoring/instructor noise — never real slide content
_NOISE_PATTERNS = [
    re.compile(r"^VO[\s«]", re.IGNORECASE),
    re.compile(r"^Slide\s+\d+$"),
    re.compile(r"OST and Script", re.IGNORECASE),
    re.compile(r"Fully Shared Slide", re.IGNORECASE),
    re.compile(r"^Course Name\s*:", re.IGNORECASE),
    re.compile(r"^Module Name\s*:", re.IGNORECASE),
    re.compile(r"The diagram is already created", re.IGNORECASE),
    re.compile(r"Please decide a presentation template", re.IGNORECASE),
    re.compile(r"^OST\s", re.IGNORECASE),
    re.compile(r"^Versal Gen\d", re.IGNORECASE),
    re.compile(r"AMD Versal.*Gen 2.*Programmable Logic", re.IGNORECASE),
    re.compile(r"^local instruction & data memory", re.IGNORECASE),
    re.compile(r"^RPU$", re.IGNORECASE),
    # Headers / footers / boilerplate
    re.compile(r"^(©|copyright|\(c\))\s", re.IGNORECASE),
    re.compile(r"all rights reserved", re.IGNORECASE),
    re.compile(r"^confidential\b", re.IGNORECASE),
    re.compile(r"^internal use only\b", re.IGNORECASE),
    re.compile(r"^(duration|learning path|level)\s*:", re.IGNORECASE),
    re.compile(r"page\s+\d+\s+of\s+\d+", re.IGNORECASE),
]


def is_noise(line):
    return any(p.search(line) for p in _NOISE_PATTERNS)


def is_strikethrough(run):
    """Return True if this run has strikethrough formatting (any OOXML variant)."""
    try:
        rPr = run._r.get_or_add_rPr()
        for tag in ("strike", "dstrike"):
            for ns in (
                "http://schemas.openxmlformats.org/drawingml/2006/main",
                "http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing",
            ):
                el = rPr.find(f".//{{{ns}}}{tag}")
                if el is not None and el.get("val", "sng") not in ("0", "false", "noStrike"):
                    return True
    except Exception:
        pass
    return False


def describe_image(shape):
    """Return the most useful description of an image shape."""
    parts = []
    try:
        nvpic = shape._element.nvPicPr
        cNvPr = nvpic.cNvPr
        title = cNvPr.get("title", "").strip()
        descr = cNvPr.get("descr", "").strip()
        if title:
            parts.append(title)
        if descr and descr != title:
            parts.append(descr)
    except Exception:
        pass
    # Fall back to shape name only if it looks meaningful (not "Picture 3")
    if not parts and shape.name and not re.match(r"^Picture\s+\d+$", shape.name, re.IGNORECASE):
        parts.append(shape.name)
    return " — ".join(parts) if parts else ""


def extract_shape_text(shape):
    """Recursively extract non-strikethrough text from a shape."""
    texts = []
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            texts.extend(extract_shape_text(s))
            # Also capture images inside groups
            if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
                desc = describe_image(s)
                if desc:
                    texts.append(f"[image: {desc}]")
    elif shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            line = " ".join(
                run.text for run in para.runs
                if run.text.strip() and not is_strikethrough(run)
            )
            if line.strip():
                texts.append(line.strip())
    return texts


def extract(pptx_path):
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"Error reading PPTX: {e}. Ensure the file is a valid, unencrypted PowerPoint.", file=sys.stderr)
        sys.exit(1)

    slides_data = []

    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        body_texts = []
        images = []

        for shape in slide.shapes:
            # Groups: recurse for text; image detection handled inside extract_shape_text
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                body_texts.extend(
                    t for t in extract_shape_text(shape) if not is_noise(t)
                )
                continue

            # Title placeholder
            try:
                if shape.placeholder_format is not None:
                    if shape.placeholder_format.idx == 0 and shape.has_text_frame:
                        title = shape.text_frame.text.strip()
                        continue
            except ValueError:
                pass

            # Text shapes
            texts = [t for t in extract_shape_text(shape) if not is_noise(t)]
            body_texts.extend(texts)

            # Top-level images (not inside groups)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                desc = describe_image(shape)
                images.append(desc if desc else "(image)")

        slides_data.append({
            "slide": i,
            "title": title,
            "text": body_texts,
            "images": images,
        })

    return slides_data


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_slides.py <path_to_pptx>", file=sys.stderr)
        sys.exit(1)

    data = extract(sys.argv[1])
    sys.stdout.buffer.write(json.dumps(data, indent=2, ensure_ascii=True).encode("utf-8"))
    sys.stdout.buffer.write(b"\n")
