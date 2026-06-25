#!/usr/bin/env python3
"""
Parse ASCII / Unicode box-drawing art into structured diagram data
(boxes + connectors), then render as native AMD-styled PowerPoint shapes.

Hybrid implementation — combines the permissive Unicode-tolerant input of the
original parser with the robust pipeline of the strict `ascii2pptx` reference
(bends, labels, line styles, emphasis, deterministic ordering, container
body rebuild).

Public API (stable — available for LLM-authored PowerPoint construction):

    parse_ascii(text)                 -> (list[Box], list[Connector])
    boxes_to_inches(boxes, conns, ...) -> (list[dict],  list[dict])
    render_to_slide(slide, boxes_in, conns_in, color_map=None)

Features supported in input:

  * Box glyphs (any mix):
      ASCII   `+ - |`
      Unicode `\u250c \u2510 \u2514 \u2518 \u251c \u2524 \u252c \u2534 \u253c \u2500 \u2501 \u2502 \u2503`
  * Line styles in connectors:
      thin    (`-` / `\u2500`)
      bold    (`=` / `\u2550`)
      dashed  (`.` / `\u254c` / `\u254d`)
  * Arrowheads:  `> < ^ v` or `\u25ba \u25c4 \u25b6 \u25c0 \u25b2 \u25bc \u2192 \u2190 \u2191 \u2193`
  * Bidirectional arrows (heads on both ends)
  * Bends via `+` junctions (after normalization) — L-shaped / Z-shaped routes
  * Inline labels: free text within 2 cells of an arrow path
  * Per-box color hints:        `# color=gold` after the top edge
  * Emphasis (accent fill):     wrap identifier as `[[ID]]`
  * Nested containers detected geometrically (unlimited depth)

Failure mode is *soft*: parse errors are logged to stderr but parse_ascii still
returns whatever boxes/connectors it could recover. Callers decide how to use or
discard partial diagrams.
"""
from __future__ import annotations

import json
import re
import sys
from collections import deque
from dataclasses import dataclass, field
from typing import List, Tuple


# ----------------------------------------------------------------------------
# Data classes (public)
# ----------------------------------------------------------------------------

@dataclass
class Box:
    id: str
    label: str
    sublabel: str = ""
    row: int = 0
    col: int = 0
    width: int = 0          # char width (corner-to-corner, inclusive)
    height: int = 0         # char height (corner-to-corner, inclusive)
    color: str = "teal"
    is_container: bool = False
    emphasized: bool = False
    children: List[str] = field(default_factory=list)


@dataclass
class Connector:
    from_id: str
    to_id: str
    direction: str = "down"        # down|up|right|left  (primary axis)
    style: str = "thin"            # thin|bold|dashed
    bidirectional: bool = False
    label: str = ""
    # Path is a list of (row, col) cells in cell-grid space, endpoints
    # included. Used by render_to_slide to draw polylines that follow
    # actual bends instead of a single straight segment.
    path: List[Tuple[int, int]] = field(default_factory=list)


# ----------------------------------------------------------------------------
# Unicode -> ASCII normalization
# ----------------------------------------------------------------------------

_UNICODE_MAP = {
    # Corners and junctions -> '+'
    "\u250c": "+", "\u2510": "+", "\u2514": "+", "\u2518": "+",
    "\u251c": "+", "\u2524": "+", "\u252c": "+", "\u2534": "+", "\u253c": "+",
    "\u2554": "+", "\u2557": "+", "\u255a": "+", "\u255d": "+",
    "\u2560": "+", "\u2563": "+", "\u2566": "+", "\u2569": "+", "\u256c": "+",
    # Horizontals -> '-'
    "\u2500": "-", "\u2501": "-", "\u2574": "-", "\u2576": "-",
    "\u257e": "-", "\u257c": "-",
    # Verticals -> '|'
    "\u2502": "|", "\u2503": "|", "\u2575": "|", "\u2577": "|",
    "\u257f": "|", "\u257d": "|",
    # Double-line horizontals -> '=' (bold style hint)
    "\u2550": "=",
    # Double-line vertical -> '|' (style determined by horizontal glyph)
    "\u2551": "|",
    # Dashed horizontals -> '.'
    "\u254c": ".", "\u254d": ".", "\u2504": ".", "\u2505": ".",
    "\u2508": ".", "\u2509": ".",
    # Dashed vertical -> '|'
    "\u254e": "|", "\u254f": "|", "\u2506": "|", "\u2507": "|",
    "\u250a": "|", "\u250b": "|",
    # Arrowheads
    "\u25ba": ">", "\u25b6": ">", "\u2192": ">", "\u27a4": ">", "\u279c": ">",
    "\u25c4": "<", "\u25c0": "<", "\u2190": "<",
    "\u25bc": "v", "\u2193": "v",
    "\u25b2": "^", "\u2191": "^",
}


def normalize(text: str) -> str:
    """Translate Unicode box-drawing chars to their ASCII equivalents."""
    return text.translate({ord(k): v for k, v in _UNICODE_MAP.items()})


# ----------------------------------------------------------------------------
# Character classes (post-normalization)
# ----------------------------------------------------------------------------

CORNER = "+"
HORIZ_THIN = "-"
HORIZ_BOLD = "="
HORIZ_DASH = "."
HORIZ_GLYPHS = set("-=.")
VERT = "|"
ARROWHEADS_H = set("<>")
ARROWHEADS_V = set("^v")
ARROWHEADS = ARROWHEADS_H | ARROWHEADS_V
LINE_GLYPHS = HORIZ_GLYPHS | {VERT, CORNER}
ARROW_GLYPHS = LINE_GLYPHS | ARROWHEADS
EDGE_OR_VERT = {CORNER, VERT}


# ----------------------------------------------------------------------------
# Grid helpers
# ----------------------------------------------------------------------------

def _load_grid(text: str) -> List[str]:
    """Normalize Unicode, expand tabs, pad rows to a common width."""
    text = normalize(text).expandtabs(4)
    lines = text.split("\n")
    while lines and lines[-1].rstrip() == "":
        lines.pop()
    width = max((len(l) for l in lines), default=0)
    return [l.ljust(width) for l in lines]


def _load_text_grid(text: str) -> List[str]:
    """Like _load_grid, but Unicode line/arrow glyphs are mapped to spaces
    instead of ASCII look-alikes. Used for label extraction so that real
    ASCII letters (v, ^, <, >) inside box content survive intact.
    """
    blank_map = {ord(k): " " for k in _UNICODE_MAP}
    text = text.translate(blank_map).expandtabs(4)
    lines = text.split("\n")
    while lines and lines[-1].rstrip() == "":
        lines.pop()
    width = max((len(l) for l in lines), default=0)
    return [l.ljust(width) for l in lines]


def _at(grid: List[str], r: int, c: int) -> str:
    if 0 <= r < len(grid) and 0 <= c < len(grid[r]):
        return grid[r][c]
    return " "


# ----------------------------------------------------------------------------
# Stage 1 - find boxes (column-drift tolerant)
# ----------------------------------------------------------------------------

def _find_boxes(grid: List[str]) -> List[Tuple[int, int, int, int]]:
    """Return inclusive (r1, c1, r2, c2) tuples for every detected box.

    Tolerant of ragged ASCII drift:
      * Corners are any `+` after normalization.
      * For each TL `+`, enumerate TR `+` candidates on the same row.
      * For each TR, enumerate BR `+` candidates *exhaustively* down the
        right column - we do NOT bail on stray glyphs along that column,
        which was the root cause of missing nested boxes in the previous
        implementation.
      * Accept if a BL `+` exists at (r2, c1) AND `_verify_edges` passes.
    """
    height = len(grid)
    width = len(grid[0]) if grid else 0
    if height == 0 or width == 0:
        return []
    plus = [[grid[r][c] == CORNER for c in range(width)] for r in range(height)]
    seen = set()
    out: List[Tuple[int, int, int, int]] = []

    for r1 in range(height - 2):
        for c1 in range(width - 2):
            if not plus[r1][c1]:
                continue
            # TL must be an OUTER corner - not a junction in a longer wall.
            # The cell above must not continue the vertical wall, and the
            # cell to the left must not continue the horizontal edge.
            if r1 > 0 and _at(grid, r1 - 1, c1) in EDGE_OR_VERT:
                continue
            if c1 > 0 and _at(grid, r1, c1 - 1) in (HORIZ_GLYPHS | {CORNER}):
                continue
            tr_cols = _scan_top_edge(grid, r1, c1)
            if not tr_cols:
                continue
            for c2 in tr_cols:
                # TR must be an outer corner too.
                if r1 > 0 and _at(grid, r1 - 1, c2) in EDGE_OR_VERT:
                    continue
                if c2 + 1 < width and _at(grid, r1, c2 + 1) in (HORIZ_GLYPHS | {CORNER}):
                    continue
                # Search for a matched (BL,BR) corner pair on any row below.
                # Uniform column drift of ±2 is allowed (body-row drift in
                # hand-drawn ASCII art is common when arrowheads sit next to
                # walls and shift downstream chars by 1-2 columns).
                found = None
                for r2 in range(r1 + 2, height):
                    for dc in (0, 1, -1, 2, -2):
                        c1b = c1 + dc
                        c2b = c2 + dc
                        if not (0 <= c1b < width and 0 <= c2b < width):
                            continue
                        if not (plus[r2][c1b] and plus[r2][c2b]):
                            continue
                        # BL and BR must be outer corners (no wall below).
                        if r2 + 1 < height:
                            if _at(grid, r2 + 1, c1b) in EDGE_OR_VERT:
                                continue
                            if _at(grid, r2 + 1, c2b) in EDGE_OR_VERT:
                                continue
                        if _verify_edges(grid, r1, c1, c2, r2, c1b, c2b):
                            found = (r2, c1b, c2b)
                            break
                    if found:
                        break
                if not found:
                    continue
                r2, c1b, c2b = found
                # Use the top-row corner columns as canonical box extent.
                # The body rows may drift by ±2, but the top edge is
                # stable and matches how a human reads the diagram.
                # Using top-row coords also keeps containment checks
                # robust: an inner box at (4,24)-(7,25) stays inside its
                # parent at (0,22)-(8,58).
                key = (r1, c1, r2, c2)
                if key in seen:
                    continue
                seen.add(key)
                out.append(key)
    return out


def _scan_top_edge(grid: List[str], r: int, c1: int) -> List[int]:
    """Walk right along row r from c1+1, listing every `+` reachable through
    HORIZ glyphs. Stops at vertical bars, arrowheads, or text.
    """
    out: List[int] = []
    width = len(grid[r]) if r < len(grid) else 0
    for c in range(c1 + 1, width):
        ch = grid[r][c]
        if ch == CORNER:
            out.append(c)
            continue
        if ch in HORIZ_GLYPHS:
            continue
        if ch == " ":
            # Tolerate spaces in the top edge (e.g. between adjacent boxes).
            continue
        # `|`, arrowhead, or text glyph - top edge ended.
        break
    return out


def _verify_edges(grid, r1, c1, c2, r2, c1b, c2b):
    """Verify a (possibly drift-shifted) rectangle is plausibly a box.

    Top corners are at (r1, c1) and (r1, c2). Bottom corners are at
    (r2, c1b) and (r2, c2b) — `c1b` and `c2b` may differ from `c1`/`c2`
    by ±2 to tolerate body-row column drift.

    * Top edge cells c1+1..c2-1 on row r1: every cell must be a HORIZ glyph
      or `+`.
    * Bottom edge cells c1b+1..c2b-1 on row r2: same rule.
    * Left vertical edge: tolerant - ≥40 % of inner-row cells at column
      c1 (or c1±1, c1±2) must be `|`/`+`.
    * Right vertical edge: same rule at column c2.
    """
    if c2 - c1 < 2 or r2 - r1 < 2 or c2b - c1b < 2:
        return False

    for c in range(c1 + 1, c2):
        ch = _at(grid, r1, c)
        if ch not in HORIZ_GLYPHS and ch != CORNER:
            return False
    for c in range(c1b + 1, c2b):
        ch = _at(grid, r2, c)
        if ch not in HORIZ_GLYPHS and ch != CORNER:
            return False

    for col in (c1, c2):
        total = r2 - r1 - 1
        good = 0
        for r in range(r1 + 1, r2):
            ok = False
            for dc in (0, -1, 1, -2, 2):
                if _at(grid, r, col + dc) in EDGE_OR_VERT:
                    ok = True
                    break
            if ok:
                good += 1
        if total > 0 and good / total < 0.4:
            return False
    return True


# ----------------------------------------------------------------------------
# Stage 2 - containment + label extraction
# ----------------------------------------------------------------------------

def _build_containment(raw_boxes):
    """Return (sorted_boxes, parent_of) where parent_of maps idx -> parent idx."""
    sorted_boxes = sorted(
        raw_boxes,
        key=lambda b: (b[2] - b[0]) * (b[3] - b[1]),
        reverse=True,
    )
    parent_of = {}
    for i, (r1, c1, r2, c2) in enumerate(sorted_boxes):
        best, best_area = None, float("inf")
        for j, (pr1, pc1, pr2, pc2) in enumerate(sorted_boxes):
            if i == j:
                continue
            # Allow ±2 column slack on the parent's left/right walls so that
            # column drift in the child's top row doesn't break containment
            # (e.g. an inner box detected at c1=63 still nests inside a
            # parent that started at c1=64).
            if pr1 < r1 and pr2 > r2 and pc1 - 2 <= c1 and pc2 + 2 >= c2:
                area = (pr2 - pr1) * (pc2 - pc1)
                if area < best_area:
                    best, best_area = j, area
        if best is not None:
            parent_of[i] = best
    return sorted_boxes, parent_of


def _is_line_glyph_in_context(text_grid, r, c, ch):
    """Decide whether an ASCII line glyph at (r,c) should be treated as a
    line-drawing char (blank it) or kept as label content (e.g. a hyphen
    inside `Pre-proc`).

    Rule: blank only if the cell is *part of a run* — i.e. an adjacent
    cell in the same orientation is also a line glyph or a wall char.
    A lone `-` or `|` surrounded by letters is preserved.
    """
    if ch in HORIZ_GLYPHS:
        left = _at(text_grid, r, c - 1)
        right = _at(text_grid, r, c + 1)
        if left in (HORIZ_GLYPHS | {CORNER, " "}) and not _is_word(left):
            if right in (HORIZ_GLYPHS | {CORNER, " "}) and not _is_word(right):
                return True
        # one side word -> hyphen
        return False
    if ch == VERT:
        up = _at(text_grid, r - 1, c)
        down = _at(text_grid, r + 1, c)
        if up in (EDGE_OR_VERT | {" "}) or down in (EDGE_OR_VERT | {" "}):
            return True
        return False
    if ch == CORNER:
        return True
    return False


def _is_word(ch):
    return ch.isalnum() or ch == "_"


def _extract_text(text_grid, r1, c1, r2, c2, child_regions):
    """Pull interior text from a box, blanking child-box footprints.

    `text_grid` has Unicode line/arrow glyphs replaced with spaces, but
    ASCII letters preserved. We additionally blank ASCII line glyphs
    (`-=.|+`) only when they appear in *line-drawing context* (adjacent
    to other line glyphs or walls), so real hyphens inside labels like
    `Pre-proc` or `Cortex-72` survive.

    Also strips trailing ` # color=NAME` comments so they don't bleed into
    the label.
    """
    def in_child(r, c):
        for cr1, cc1, cr2, cc2 in child_regions:
            if cr1 <= r <= cr2 and cc1 <= c <= cc2:
                return True
        return False

    out = []
    for r in range(r1 + 1, r2):
        buf = []
        # Read INCLUSIVE of the wall columns - walls themselves are blanked
        # in text_grid (Unicode) or by the in-context line-glyph check
        # (ASCII), so any text that drifted onto a wall column survives.
        for c in range(c1, c2 + 1):
            if in_child(r, c):
                buf.append(" ")
                continue
            ch = _at(text_grid, r, c)
            if ch in LINE_GLYPHS and _is_line_glyph_in_context(text_grid, r, c, ch):
                buf.append(" ")
            else:
                buf.append(ch)
        s = "".join(buf)
        # Strip color-hint comments before further processing.
        s = _COLOR_RE.sub("", s)
        s = re.sub(r"#\s*$", "", s)  # bare trailing '#'
        s = s.strip()
        if s:
            out.append(s)
    return out


_COLOR_RE = re.compile(r"#\s*color\s*=\s*([A-Za-z_][A-Za-z0-9_]*)")


def _extract_color_hint(grid, r1, c1, r2, c2):
    """Look for `# color=NAME` on the top-edge tail or inside the box."""
    if r1 < len(grid):
        tail = grid[r1][c2 + 1:]
        m = _COLOR_RE.search(tail)
        if m:
            return m.group(1).lower()
    for r in range(r1, r2 + 1):
        if r < len(grid):
            m = _COLOR_RE.search(grid[r])
            if m:
                return m.group(1).lower()
    return None


def _slug(s):
    s = re.sub(r"[^a-zA-Z0-9]+", "_", s.lower()).strip("_")
    return s or "box"


def _make_boxes(grid, text_grid, raw_boxes):
    sorted_boxes, parent_of = _build_containment(raw_boxes)
    boxes: List[Box] = []
    idx_to_id = {}

    children_of_i = {i: [] for i in range(len(sorted_boxes))}
    for ci, pi in parent_of.items():
        children_of_i[pi].append(ci)

    for i, (r1, c1, r2, c2) in enumerate(sorted_boxes):
        child_regions = [sorted_boxes[ci] for ci in children_of_i[i]]
        text_lines = _extract_text(text_grid, r1, c1, r2, c2, child_regions)
        label = text_lines[0] if text_lines else f"box_{i}"
        sublabel = " ".join(text_lines[1:]) if len(text_lines) > 1 else ""

        emphasized = False
        m = re.match(r"^\s*\[\[\s*(.+?)\s*\]\]\s*$", label)
        if m:
            emphasized = True
            label = m.group(1)

        label = _COLOR_RE.sub("", label).strip()
        sublabel = _COLOR_RE.sub("", sublabel).strip()

        box_id = _slug(label)
        orig_id, n = box_id, 2
        while any(b.id == box_id for b in boxes):
            box_id = f"{orig_id}_{n}"
            n += 1

        color = _extract_color_hint(grid, r1, c1, r2, c2) or "teal"
        is_container = bool(children_of_i[i])
        if is_container and color == "teal":
            color = "dark"

        boxes.append(Box(
            id=box_id, label=label, sublabel=sublabel,
            row=r1, col=c1,
            width=c2 - c1, height=r2 - r1,
            color=color, is_container=is_container,
            emphasized=emphasized,
        ))
        idx_to_id[i] = box_id

    for ci, pi in parent_of.items():
        if ci in idx_to_id and pi in idx_to_id:
            parent_box = next(b for b in boxes if b.id == idx_to_id[pi])
            parent_box.children.append(idx_to_id[ci])

    boxes.sort(key=lambda b: (b.row, b.col, b.id))
    return boxes


# ----------------------------------------------------------------------------
# Stage 3 - find connectors via BFS components
# ----------------------------------------------------------------------------

def _border_cells(b: Box):
    r1, c1, r2, c2 = b.row, b.col, b.row + b.height, b.col + b.width
    cells = set()
    for c in range(c1, c2 + 1):
        cells.add((r1, c))
        cells.add((r2, c))
    for r in range(r1 + 1, r2):
        cells.add((r, c1))
        cells.add((r, c2))
    return cells


def _build_border_map(boxes):
    out = {}
    for b in boxes:
        for cell in _border_cells(b):
            out.setdefault(cell, []).append(b)
    return out


def _find_connectors(grid, boxes):
    """Find connected components of arrow glyphs outside box borders."""
    height = len(grid)
    width = len(grid[0]) if grid else 0
    border_map = _build_border_map(boxes)
    border_cells = set(border_map.keys())
    visited = [[False] * width for _ in range(height)]
    connectors: List[Connector] = []

    for r in range(height):
        for c in range(width):
            if visited[r][c]:
                continue
            ch = grid[r][c]
            if ch not in ARROW_GLYPHS:
                continue
            if (r, c) in border_cells:
                visited[r][c] = True
                continue
            comp = _flood(grid, r, c, visited, border_cells)
            if len(comp) < 2:
                continue
            conn = _component_to_connector(grid, comp, border_map)
            if conn is not None:
                connectors.append(conn)

    uniq = []
    seen = set()
    for c in connectors:
        key = (c.from_id, c.to_id, c.style, c.bidirectional)
        rev = (c.to_id, c.from_id, c.style, c.bidirectional)
        if key in seen or (c.bidirectional and rev in seen):
            continue
        seen.add(key)
        uniq.append(c)
    uniq.sort(key=lambda c: (c.from_id, c.to_id))
    return uniq


def _flood(grid, r0, c0, visited, border_cells):
    height = len(grid)
    width = len(grid[0]) if grid else 0
    comp = []
    q = deque([(r0, c0)])
    visited[r0][c0] = True
    while q:
        r, c = q.popleft()
        comp.append((r, c))
        for nr, nc in ((r + 1, c), (r - 1, c), (r, c + 1), (r, c - 1)):
            if not (0 <= nr < height and 0 <= nc < width):
                continue
            if visited[nr][nc]:
                continue
            if grid[nr][nc] not in ARROW_GLYPHS:
                continue
            if (nr, nc) in border_cells:
                continue
            visited[nr][nc] = True
            q.append((nr, nc))
    return comp


def _component_to_connector(grid, comp, border_map):
    cells = set(comp)
    deg = {}
    for (r, c) in comp:
        d = 0
        for nr, nc in ((r + 1, c), (r - 1, c), (r, c + 1), (r, c - 1)):
            if (nr, nc) in cells:
                d += 1
        deg[(r, c)] = d
    endpoints = [cell for cell, d in deg.items() if d <= 1]
    if len(endpoints) < 2:
        return None
    if len(endpoints) > 2:
        anchored = [ep for ep in endpoints if _anchor_box(ep, border_map, grid)]
        if len(anchored) >= 2:
            endpoints = anchored
    if len(endpoints) < 2:
        return None

    pairs = []
    for i in range(len(endpoints)):
        for j in range(i + 1, len(endpoints)):
            a, b = endpoints[i], endpoints[j]
            ba = _anchor_box(a, border_map, grid)
            bb = _anchor_box(b, border_map, grid)
            if ba is not None and bb is not None and ba.id != bb.id:
                pairs.append((a, b, ba, bb))
    if not pairs:
        return None
    pairs.sort(key=lambda p: (p[0][0] + p[0][1] + p[1][0] + p[1][1]))
    ep_a, ep_b, box_a, box_b = pairs[0]

    head_a = grid[ep_a[0]][ep_a[1]] in ARROWHEADS
    head_b = grid[ep_b[0]][ep_b[1]] in ARROWHEADS

    if head_a and head_b:
        bidir = True
        if box_a.id < box_b.id:
            src_box, dst_box, src_ep, dst_ep = box_a, box_b, ep_a, ep_b
        else:
            src_box, dst_box, src_ep, dst_ep = box_b, box_a, ep_b, ep_a
    elif head_b:
        bidir = False
        src_box, dst_box, src_ep, dst_ep = box_a, box_b, ep_a, ep_b
    elif head_a:
        bidir = False
        src_box, dst_box, src_ep, dst_ep = box_b, box_a, ep_b, ep_a
    else:
        bidir = False
        if (ep_a[0], ep_a[1]) <= (ep_b[0], ep_b[1]):
            src_box, dst_box, src_ep, dst_ep = box_a, box_b, ep_a, ep_b
        else:
            src_box, dst_box, src_ep, dst_ep = box_b, box_a, ep_b, ep_a

    path = _walk_path(cells, src_ep, dst_ep)
    style = _classify_style(grid, path)
    dr = dst_ep[0] - src_ep[0]
    dc = dst_ep[1] - src_ep[1]
    if abs(dr) >= abs(dc):
        direction = "down" if dr > 0 else "up"
    else:
        direction = "right" if dc > 0 else "left"

    return Connector(
        from_id=src_box.id, to_id=dst_box.id,
        direction=direction, style=style,
        bidirectional=bidir, label="", path=path,
    )


def _anchor_box(cell, border_map, grid=None):
    """Find the box whose border this connector endpoint touches.

    Tolerant of small gaps (up to 2 spaces) between an arrowhead and the
    target box wall - common when authors write `══►  │` with a space
    before the wall.

    If ``grid`` is provided and the cell is a directional arrowhead, we
    extend the search in the arrow's pointing direction up to 2 cells
    beyond the immediate neighbour.
    """
    r, c = cell
    # Direct neighbours
    for nr, nc in ((r + 1, c), (r - 1, c), (r, c + 1), (r, c - 1)):
        if (nr, nc) in border_map:
            return border_map[(nr, nc)][0]

    if grid is None:
        return None
    ch = grid[r][c] if 0 <= r < len(grid) and 0 <= c < len(grid[r]) else " "
    # Directional extension for arrowheads
    dirs = []
    if ch == ">":
        dirs.append((0, 1))
    elif ch == "<":
        dirs.append((0, -1))
    elif ch == "v":
        dirs.append((1, 0))
    elif ch == "^":
        dirs.append((-1, 0))
    else:
        # Tail end - try all four directions for a 2-cell gap.
        dirs = [(0, 1), (0, -1), (1, 0), (-1, 0)]
    for dr, dc in dirs:
        for k in (2, 3):
            nr, nc = r + dr * k, c + dc * k
            if (nr, nc) in border_map:
                return border_map[(nr, nc)][0]
            # Stop scanning if a non-space, non-arrow cell blocks the path.
            if 0 <= nr < len(grid) and 0 <= nc < len(grid[nr]):
                blocker = grid[nr][nc]
                if blocker != " " and blocker not in ARROW_GLYPHS:
                    break
    return None


def _walk_path(cells, start, end):
    prev = {start: None}
    q = deque([start])
    while q:
        cur = q.popleft()
        if cur == end:
            break
        r, c = cur
        for nr, nc in ((r + 1, c), (r - 1, c), (r, c + 1), (r, c - 1)):
            if (nr, nc) in cells and (nr, nc) not in prev:
                prev[(nr, nc)] = cur
                q.append((nr, nc))
    if end not in prev:
        return [start, end]
    out = []
    cur = end
    while cur is not None:
        out.append(cur)
        cur = prev[cur]
    out.reverse()
    return out


def _classify_style(grid, path):
    bold = dash = False
    for (r, c) in path:
        if 0 <= r < len(grid) and 0 <= c < len(grid[r]):
            ch = grid[r][c]
            if ch == HORIZ_BOLD:
                bold = True
            elif ch == HORIZ_DASH:
                dash = True
    if bold:
        return "bold"
    if dash:
        return "dashed"
    return "thin"


# ----------------------------------------------------------------------------
# Stage 4 - attach edge labels
# ----------------------------------------------------------------------------

def _attach_labels(grid, boxes, connectors):
    if not connectors:
        return
    height = len(grid)
    width = len(grid[0]) if grid else 0

    border_cells = set()
    for b in boxes:
        border_cells |= _border_cells(b)
    interior = set()
    for b in boxes:
        for r in range(b.row + 1, b.row + b.height):
            for c in range(b.col + 1, b.col + b.width):
                interior.add((r, c))
    arrow_cells = set()
    for conn in connectors:
        arrow_cells.update(conn.path)

    for r in range(height):
        c = 0
        while c < width:
            ch = grid[r][c]
            if (ch == " " or (r, c) in border_cells or (r, c) in interior
                    or (r, c) in arrow_cells or ch in ARROW_GLYPHS):
                c += 1
                continue
            start = c
            buf = []
            while c < width:
                ch2 = grid[r][c]
                if (ch2 == " " or (r, c) in border_cells or (r, c) in interior
                        or (r, c) in arrow_cells or ch2 in ARROW_GLYPHS):
                    break
                buf.append(ch2)
                c += 1
            text = "".join(buf).strip()
            if not text or text.startswith("#"):
                continue
            best = None
            best_d = 10**9
            for ci, conn in enumerate(connectors):
                d = _min_distance(start, r, len(buf), conn.path)
                if d < best_d:
                    best_d = d
                    best = ci
            if best is not None and best_d <= 2 and not connectors[best].label:
                connectors[best].label = text


def _min_distance(rx, ry, rlen, path):
    best = 10**9
    for (pr, pc) in path:
        for k in range(rlen):
            d = abs(pr - ry) + abs(pc - (rx + k))
            if d < best:
                best = d
    return best


# ----------------------------------------------------------------------------
# Public parser
# ----------------------------------------------------------------------------

def parse_ascii(text: str):
    """Parse ASCII / Unicode box-drawing art -> (list[Box], list[Connector]).

    Errors are not raised: a warning is printed to stderr and partial
    results are returned. If no boxes can be recovered, returns ([], []) so
    the caller can fall back to a non-diagram layout.
    """
    try:
        grid = _load_grid(text)
        text_grid = _load_text_grid(text)
        if not grid:
            return [], []
        raw_boxes = _find_boxes(grid)
        if not raw_boxes:
            return [], []
        boxes = _make_boxes(grid, text_grid, raw_boxes)
        connectors = _find_connectors(grid, boxes)
        _attach_labels(grid, boxes, connectors)
        return boxes, connectors
    except Exception as exc:  # pragma: no cover
        print(f"ascii_to_diagram: parse error: {exc}", file=sys.stderr)
        return [], []


# ----------------------------------------------------------------------------
# Coordinate mapping
# ----------------------------------------------------------------------------

def boxes_to_inches(boxes, connectors,
                    slide_left=0.4, slide_top=1.2,
                    slide_width=12.0, slide_height=5.8):
    """Map char-grid to inch coordinates with aspect preservation.

    Each cell is treated as ~2:1 (twice as wide as tall) which matches how
    monospace fonts render box-drawing art. The diagram is centered.
    """
    if not boxes:
        return [], []

    min_r = min(b.row for b in boxes)
    min_c = min(b.col for b in boxes)
    max_r = max(b.row + b.height for b in boxes)
    max_c = max(b.col + b.width for b in boxes)
    char_h = max(max_r - min_r, 1)
    char_w = max(max_c - min_c, 1)

    pad = 0.94
    avail_w = slide_width * pad
    avail_h = slide_height * pad
    sx_max = avail_w / char_w
    sy_max = avail_h / char_h
    sx = min(sx_max, sy_max * 2)
    sy = sx / 2

    used_w = char_w * sx
    used_h = char_h * sy
    ox = slide_left + (slide_width - used_w) / 2
    oy = slide_top + (slide_height - used_h) / 2

    boxes_out = []
    for b in boxes:
        x = ox + (b.col - min_c) * sx
        y = oy + (b.row - min_r) * sy
        w = max(b.width * sx, 0.8)
        h = max(b.height * sy, 0.35)
        boxes_out.append({
            "id": b.id,
            "label": b.label,
            "sublabel": b.sublabel,
            "x": round(x, 3),
            "y": round(y, 3),
            "w": round(w, 3),
            "h": round(h, 3),
            "color": b.color,
            "is_container": b.is_container,
            "emphasized": b.emphasized,
            "children": list(b.children),
        })

    def cell_to_inch(r, c):
        x = ox + (c - min_c) * sx + sx / 2
        y = oy + (r - min_r) * sy + sy / 2
        return round(x, 3), round(y, 3)

    conns_out = []
    for c in connectors:
        path_in = [cell_to_inch(r, cc) for (r, cc) in c.path] if c.path else []
        conns_out.append({
            "from": c.from_id,
            "to": c.to_id,
            "direction": c.direction,
            "style": c.style,
            "bidirectional": c.bidirectional,
            "label": c.label,
            "path": path_in,
        })
    return boxes_out, conns_out


# ----------------------------------------------------------------------------
# PowerPoint rendering (AMD palette)
# ----------------------------------------------------------------------------

def render_to_slide(slide, boxes_inches, connectors_inches, color_map=None):
    """Render parsed diagram onto a python-pptx Slide using AMD styling."""
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    PALETTE = {
        "teal":  RGBColor(0,   194, 222),
        "gold":  RGBColor(193, 169, 104),
        "red":   RGBColor(237, 28,  36),
        "green": RGBColor(0,   178, 86),
        "blue":  RGBColor(58,  119, 195),
        "dark":  RGBColor(38,  38,  38),
        "mid":   RGBColor(64,  64,  64),
        "black": RGBColor(0,   0,   0),
        "white": RGBColor(200, 200, 200),
    }
    if color_map:
        PALETTE.update(color_map)

    WHITE = RGBColor(255, 255, 255)
    LIGHT = RGBColor(225, 225, 225)
    BORDER = RGBColor(64, 64, 64)
    EMPH_BORDER = RGBColor(255, 198, 0)
    LINE_COLOR = RGBColor(0, 194, 222)

    shape_map = {}

    box_order = sorted(boxes_inches,
                       key=lambda b: 0 if b.get("is_container") else 1)
    for box in box_order:
        x, y, w, h = box["x"], box["y"], box["w"], box["h"]
        bg = PALETTE.get(box.get("color", "teal"), PALETTE["teal"])

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
        if box.get("emphasized"):
            shape.line.color.rgb = EMPH_BORDER
            shape.line.width = Pt(2.5)
        else:
            shape.line.color.rgb = BORDER
            shape.line.width = Pt(1)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.06)
        tf.margin_right = Inches(0.06)
        tf.margin_top = Inches(0.04)
        tf.margin_bottom = Inches(0.04)
        tf.clear()

        p = tf.paragraphs[0]
        p.text = box["label"]
        p.font.name = "Arial"
        p.font.bold = True
        p.font.color.rgb = WHITE
        if box.get("is_container"):
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.LEFT
        else:
            p.font.size = Pt(13)
            p.alignment = PP_ALIGN.CENTER

        if box.get("sublabel"):
            p2 = tf.add_paragraph()
            p2.text = box["sublabel"]
            p2.font.name = "Arial"
            p2.font.size = Pt(10)
            p2.font.color.rgb = LIGHT
            p2.alignment = (PP_ALIGN.LEFT if box.get("is_container")
                            else PP_ALIGN.CENTER)

        shape_map[box["id"]] = shape

    for conn in connectors_inches:
        src = shape_map.get(conn["from"])
        dst = shape_map.get(conn["to"])
        if not src or not dst:
            continue

        path = conn.get("path") or []
        bidir = bool(conn.get("bidirectional"))
        style = conn.get("style", "thin")

        if path:
            start_pt = _snap_to_box_edge(_inch_to_emu(path[0]), src)
            end_pt = _snap_to_box_edge(_inch_to_emu(path[-1]), dst)
            mid_pts = [_inch_to_emu(p) for p in path[1:-1]]
            pts = [start_pt] + mid_pts + [end_pt]
            bends = _bend_points(pts)
        else:
            bends = _orthogonal_route(src, dst,
                                      conn.get("direction", "right"))

        _draw_polyline(slide, bends, style, bidir, LINE_COLOR)

        if conn.get("label"):
            _draw_edge_label(slide, conn["label"], bends)


def _inch_to_emu(pt_inches):
    from pptx.util import Inches
    x, y = pt_inches
    return Inches(x), Inches(y)


def _snap_to_box_edge(pt, shape):
    px, py = pt
    sx, sy = shape.left, shape.top
    sw, sh = shape.width, shape.height
    candidates = [
        (sx, max(sy, min(py, sy + sh))),
        (sx + sw, max(sy, min(py, sy + sh))),
        (max(sx, min(px, sx + sw)), sy),
        (max(sx, min(px, sx + sw)), sy + sh),
    ]
    candidates.sort(key=lambda c: (c[0] - px) ** 2 + (c[1] - py) ** 2)
    return candidates[0]


def _orthogonal_route(src, dst, direction):
    sl, st, sw, sh = src.left, src.top, src.width, src.height
    dl, dt, dw, dh = dst.left, dst.top, dst.width, dst.height
    if direction == "down":
        return [(sl + sw // 2, st + sh), (dl + dw // 2, dt)]
    if direction == "up":
        return [(sl + sw // 2, st), (dl + dw // 2, dt + dh)]
    if direction == "right":
        return [(sl + sw, st + sh // 2), (dl, dt + dh // 2)]
    return [(sl, st + sh // 2), (dl + dw, dt + dh // 2)]


def _bend_points(pts):
    if len(pts) <= 2:
        return list(pts)
    out = [pts[0]]
    for i in range(1, len(pts) - 1):
        ax, ay = pts[i - 1]
        bx, by = pts[i]
        cx, cy = pts[i + 1]
        dxa = _sgn(bx - ax); dya = _sgn(by - ay)
        dxb = _sgn(cx - bx); dyb = _sgn(cy - by)
        if (dxa, dya) != (dxb, dyb):
            out.append(pts[i])
    out.append(pts[-1])
    return out


def _sgn(v):
    return (v > 0) - (v < 0)


def _draw_polyline(slide, pts, style, bidir, color):
    from lxml import etree
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.util import Pt

    A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

    if len(pts) < 2:
        return

    if len(pts) == 2:
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, pts[0][0], pts[0][1], pts[1][0], pts[1][1],
        )
        line = conn.line
    else:
        try:
            builder = slide.shapes.build_freeform(pts[0][0], pts[0][1],
                                                  scale=1.0)
            builder.add_line_segments([(p[0], p[1]) for p in pts[1:]],
                                      close=False)
            shape = builder.convert_to_shape()
            shape.fill.background()
            line = shape.line
        except Exception:
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, pts[0][0], pts[0][1],
                pts[-1][0], pts[-1][1],
            )
            line = conn.line

    line.color.rgb = color
    if style == "bold":
        line.width = Pt(2.5)
    elif style == "dashed":
        line.width = Pt(1.5)
        line.dash_style = MSO_LINE_DASH_STYLE.DASH
    else:
        line.width = Pt(1.5)

    try:
        ln = line._get_or_add_ln()
        for tag in ("headEnd", "tailEnd"):
            for el in ln.findall(f"{{{A_NS}}}{tag}"):
                ln.remove(el)
        tail = etree.SubElement(
            ln, f"{{{A_NS}}}tailEnd",
            attrib={"type": "triangle", "w": "med", "len": "med"},
        )
        if bidir:
            head = etree.SubElement(
                ln, f"{{{A_NS}}}headEnd",
                attrib={"type": "triangle", "w": "med", "len": "med"},
            )
            ln.remove(head)
            ln.insert(list(ln).index(tail), head)
    except Exception:
        pass


def _draw_edge_label(slide, text, pts):
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    if not pts:
        return
    mid = len(pts) // 2
    cx, cy = pts[mid]
    w = max(Inches(0.6), Inches(len(text) * 0.08))
    h = Inches(0.25)
    tb = slide.shapes.add_textbox(cx - w // 2, cy - h // 2, w, h)
    tb.fill.solid()
    tb.fill.fore_color.rgb = RGBColor(38, 38, 38)
    tb.line.fill.background()
    tf = tb.text_frame
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(225, 225, 225)


# ----------------------------------------------------------------------------
# CLI (debug helper)
# ----------------------------------------------------------------------------

def main():
    import argparse
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--input", required=True, help="ASCII art text file")
    ap.add_argument("--output", required=True,
                    help="Output JSON with boxes/connectors (debug dump)")
    args = ap.parse_args()

    text = open(args.input, encoding="utf-8").read()
    boxes, connectors = parse_ascii(text)
    boxes_in, conns_in = boxes_to_inches(boxes, connectors)
    result = {"boxes": boxes_in, "connectors": conns_in}
    with open(args.output, "w", encoding="utf-8") as f:
        json.dump(result, f, indent=2)
    print(f"Parsed {len(boxes_in)} boxes, {len(conns_in)} connectors "
          f"-> {args.output}")


if __name__ == "__main__":
    main()
