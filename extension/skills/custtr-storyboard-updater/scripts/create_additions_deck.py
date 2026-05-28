#!/usr/bin/env python3
"""
Generate a standalone additions PPTX from plan actions using python-pptx.

Supports multiple slide layouts specified in each action's "slide_layout" field:
  - comparison_table: native PPTX table with header row and alternating shading
  - block_diagram:    boxes + connectors showing architecture relationships
  - two_column:       left/right content areas (before/after, concept/detail)
  - cards:            2-4 stacked or side-by-side info cards (default fallback)
  - key_takeaway:     single bold statement centered on slide

If "slide_layout" is omitted, defaults to "cards" for backwards compatibility.
"""
import argparse
import json
import sys
from pathlib import Path

# Ensure sibling scripts are importable regardless of working directory
sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

from ascii_to_diagram import parse_ascii, boxes_to_inches, render_to_slide
from json_helpers import safe_load_json  # noqa: E402

# ── Brand colours ──────────────────────────────────────────────
BLACK   = RGBColor(0, 0, 0)
WHITE   = RGBColor(255, 255, 255)
TEAL    = RGBColor(0, 194, 222)
GOLD    = RGBColor(193, 169, 104)
RED     = RGBColor(237, 28, 36)
DARK    = RGBColor(22, 22, 22)
DARK2   = RGBColor(38, 38, 38)
MID     = RGBColor(64, 64, 64)
LIGHT   = RGBColor(157, 159, 162)
HDR_BG  = RGBColor(0, 60, 80)

SLIDE_W = 12192000
SLIDE_H = 6858000


# ── Helpers ────────────────────────────────────────────────────
def set_font(run, size, color=WHITE, bold=False):
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def para(tf, text, size=14, color=WHITE, bold=False, alignment=None):
    p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
    p.text = str(text)
    if p.runs:
        set_font(p.runs[0], size, color, bold)
    p.space_after = Pt(4)
    if alignment:
        p.alignment = alignment
    return p


def fill_shape(shape, color, line=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()


def add_badge(slide):
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(11.55), Inches(0.22), Inches(1.18), Inches(0.32),
    )
    fill_shape(badge, RGBColor(255, 255, 0))
    tf = badge.text_frame
    tf.clear()
    p = para(tf, "New Slide", 10, RGBColor(0, 0, 0), True, PP_ALIGN.CENTER)


def add_title_bar(slide, title):
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.45), Inches(0.48), Inches(0.06), Inches(0.48),
    )
    fill_shape(accent, TEAL)
    box = slide.shapes.add_textbox(
        Inches(0.62), Inches(0.36), Inches(10.7), Inches(0.72),
    )
    tf = box.text_frame
    tf.clear()
    para(tf, title, 28, WHITE, True)
    add_badge(slide)


def add_notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text


def make_slide(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BLACK
    return slide


# ── Validation ─────────────────────────────────────────────────
def validate_required_fields(action):
    """Check that all mandatory metadata fields are present and speaker_notes is long enough."""
    required = ["title", "learning_goal", "why_this_slide_exists",
                 "what_customer_should_understand", "speaker_notes"]
    missing = [k for k in required if not action.get(k)]
    if missing:
        raise ValueError(
            f"add_new_slide '{action.get('title', '<untitled>')}' missing: {missing}"
        )
    words = len(str(action.get("speaker_notes", "")).split())
    if words < 80:
        raise ValueError(
            f"add_new_slide '{action.get('title')}' needs >=80 words in speaker_notes; got {words}"
        )


def validate_layout_schema(action):
    """Validate slide_layout value and its required data field."""
    layout = action.get("slide_layout", "")
    if not layout:
        raise ValueError(
            f"add_new_slide '{action.get('title', '<untitled>')}' is missing "
            f"slide_layout; must be one of: {sorted(LAYOUTS.keys())}"
        )
    if layout not in LAYOUTS:
        raise ValueError(
            f"add_new_slide '{action.get('title', '<untitled>')}' uses unknown "
            f"slide_layout '{layout}'; valid layouts: {sorted(LAYOUTS.keys())}"
        )
    _LAYOUT_DATA_FIELDS = {
        "comparison_table": "table", "block_diagram": "diagram",
        "ascii_diagram": "ascii_art", "two_column": "columns",
        "key_takeaway": "statement", "cards": "cards",
    }
    data_field = _LAYOUT_DATA_FIELDS.get(layout)
    if data_field and not action.get(data_field):
        raise ValueError(
            f"add_new_slide '{action.get('title', '<untitled>')}' uses layout "
            f"'{layout}' but is missing required field '{data_field}'"
        )
    if action.get("body"):
        raise ValueError(
            f"add_new_slide '{action.get('title', '<untitled>')}': 'body' is not "
            f"a valid field for any layout; use the layout-specific data field "
            f"(e.g., 'cards', 'table', 'columns', 'diagram', 'ascii_art', 'statement')"
        )


def validate_content_quality(action):
    """Strict quality checks — diagram/table source attribution, card headings, column bullets."""
    layout = action.get("slide_layout", "cards")
    diagram_layouts = ("block_diagram", "ascii_diagram")
    if layout in diagram_layouts:
        if layout == "block_diagram":
            boxes = (action.get("diagram") or {}).get("boxes") or []
            if not boxes:
                raise ValueError(
                    f"add_new_slide '{action.get('title')}' uses block_diagram "
                    f"but has no diagram.boxes; quality bar requires "
                    f">=1 grouped region / functional block."
                )
        elif layout == "ascii_diagram":
            if not (action.get("ascii_art") or "").strip():
                raise ValueError(
                    f"add_new_slide '{action.get('title')}' uses ascii_diagram "
                    f"with empty ascii_art."
                )
        if not (action.get("source_footer") or "").strip() and not (action.get("source_basis") or "").strip():
            raise ValueError(
                f"add_new_slide '{action.get('title')}' is a diagram slide but "
                f"has no source_footer or source_basis; quality bar requires "
                f"explicit source attribution (e.g. 'Source: PG346 v3.x §4')."
            )
        kp = action.get("key_points")
        if kp and isinstance(kp, dict):
            kp_bullets = kp.get("bullets", [])
            if len(kp_bullets) < 3:
                import warnings
                warnings.warn(
                    f"add_new_slide '{action.get('title')}': key_points has "
                    f"{len(kp_bullets)} bullets (recommend 3-6 for visible OST)."
                )
            notes = (action.get("speaker_notes") or "").lower()
            if notes and kp_bullets:
                for bi, bullet in enumerate(kp_bullets):
                    words = bullet.split()
                    key_term = words[0].rstrip(":,;") if words else ""
                    if len(key_term) < 3 and len(words) > 1:
                        key_term = " ".join(words[:2]).rstrip(":,;")
                    if key_term and key_term.lower() not in notes:
                        raise ValueError(
                            f"add_new_slide '{action.get('title')}': key_points "
                            f"bullet {bi} key term '{key_term}' not found in "
                            f"speaker_notes — OST and narration must stay in sync "
                            f"(Rule 34)."
                        )
        elif not kp:
            import warnings
            warnings.warn(
                f"add_new_slide '{action.get('title')}': diagram layout without "
                f"key_points — consider adding 3-6 visible bullet points for "
                f"narrated storyboards."
            )
    if layout == "comparison_table":
        tbl = action.get("table") or {}
        headers = tbl.get("headers") or []
        rows = tbl.get("rows") or []
        if len(headers) < 3:
            raise ValueError(
                f"add_new_slide '{action.get('title')}' comparison_table needs "
                f">=3 columns (got {len(headers)})."
            )
        if len(rows) < 2:
            raise ValueError(
                f"add_new_slide '{action.get('title')}' comparison_table needs "
                f">=2 rows of concrete values (got {len(rows)})."
            )
        if not (action.get("source_footer") or "").strip() and not (action.get("source_basis") or "").strip():
            raise ValueError(
                f"add_new_slide '{action.get('title')}' is a comparison_table slide "
                f"but has no source_footer or source_basis."
            )
    if layout == "cards":
        card_list = action.get("cards") or []
        if len(card_list) < 2:
            raise ValueError(
                f"add_new_slide '{action.get('title')}' cards layout needs "
                f">=2 cards (got {len(card_list)})."
            )
        _generic = {"what it is", "why it matters", "customer value",
                    "what", "why", "value", "overview", "benefits"}
        for ci, card in enumerate(card_list):
            heading = (card.get("heading") or "").strip().lower()
            if heading in _generic:
                raise ValueError(
                    f"add_new_slide '{action.get('title')}' cards[{ci}] heading "
                    f"'{card.get('heading')}' is a generic anti-pattern; use "
                    f"domain-specific headings (component names, feature names)."
                )
    if layout == "two_column":
        cols = action.get("columns") or []
        if len(cols) < 2:
            raise ValueError(
                f"add_new_slide '{action.get('title')}' two_column needs "
                f">=2 columns (got {len(cols)})."
            )
        for ci, col in enumerate(cols[:2]):
            if not col.get("bullets"):
                raise ValueError(
                    f"add_new_slide '{action.get('title')}' columns[{ci}] "
                    f"has no bullets."
                )
    if layout == "key_takeaway":
        if not (action.get("statement") or "").strip():
            raise ValueError(
                f"add_new_slide '{action.get('title')}' key_takeaway needs "
                f"a non-empty 'statement' field."
            )


def validate_new_slide_action(action, strict_quality: bool = False):
    validate_required_fields(action)
    validate_layout_schema(action)
    if strict_quality:
        validate_content_quality(action)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# LAYOUT: comparison_table
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def layout_comparison_table(slide, action):
    """Native PPTX table with header row and alternating row shading.

    Plan action fields:
      table.headers: ["Criteria", "APU", "RPU"]
      table.rows: [["Clock", "1.7 GHz", "600 MHz"], ...]
    """
    table_spec = action.get("table", {})
    headers = table_spec.get("headers", ["Criteria", "Option A", "Option B"])
    rows_data = table_spec.get("rows", [])
    if not rows_data:
        rows_data = [["(no data)", "", ""]]

    n_cols = len(headers)
    n_rows = len(rows_data) + 1

    left = Inches(0.5)
    top = Inches(1.25)
    width = Inches(12.3)
    total_h = min(Inches(5.8), Inches(0.50) * n_rows)

    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, total_h)
    table = tbl_shape.table

    first_w = int(width * 0.22)
    rest_w = int((width - first_w) / max(n_cols - 1, 1))
    table.columns[0].width = Emu(first_w)
    for c in range(1, n_cols):
        table.columns[c].width = Emu(rest_w)

    for c, hdr in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = hdr
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Arial"
        cell.fill.solid()
        cell.fill.fore_color.rgb = HDR_BG

    for r, row in enumerate(rows_data):
        for c in range(n_cols):
            cell = table.cell(r + 1, c)
            cell.text = row[c] if c < len(row) else ""
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = WHITE
                p.font.name = "Arial"
                if c == 0:
                    p.font.bold = True
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK if r % 2 == 0 else DARK2


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# LAYOUT: block_diagram
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def layout_block_diagram(slide, action):
    """Render boxes and connectors from a structured spec.

    Plan action fields:
      diagram.boxes: [{"id":"apu", "label":"APU", "sublabel":"Cortex-A72",
                        "x":1, "y":2, "w":2.5, "h":1.0, "color":"teal"}, ...]
      diagram.connectors: [{"from":"apu", "to":"interconnect"}, ...]

    Supported colors: teal, gold, red, dark, mid, black (defaults to teal).
    """
    diagram = action.get("diagram", {})
    boxes = diagram.get("boxes", [])
    connectors = diagram.get("connectors", [])

    color_map = {
        "teal": TEAL, "gold": GOLD, "red": RED,
        "dark": DARK, "mid": MID, "black": BLACK,
    }

    box_shapes = {}
    for box in boxes:
        x = Inches(box.get("x", 1))
        y = Inches(box.get("y", 2))
        w = Inches(box.get("w", 2.5))
        h = Inches(box.get("h", 1.0))
        bg = color_map.get(box.get("color", "teal"), TEAL)

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        fill_shape(shape, bg, MID)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.clear()
        para(tf, box.get("label", ""), 14, WHITE, True, PP_ALIGN.CENTER)
        if box.get("sublabel"):
            para(tf, box["sublabel"], 11, LIGHT, False, PP_ALIGN.CENTER)

        box_shapes[box.get("id", box.get("label", ""))] = shape

    for conn in connectors:
        src = box_shapes.get(conn.get("from"))
        dst = box_shapes.get(conn.get("to"))
        if src and dst:
            x1 = src.left + src.width // 2
            y1 = src.top + src.height
            x2 = dst.left + dst.width // 2
            y2 = dst.top
            line = slide.shapes.add_connector(1, x1, y1, x2, y2)
            line.line.color.rgb = TEAL
            line.line.width = Pt(2)

    if action.get("key_points"):
        _add_key_points_panel(slide, action["key_points"])


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# LAYOUT: two_column
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def layout_two_column(slide, action):
    """Two side-by-side content areas with headings.

    Plan action fields:
      columns: [
        {"heading":"APU", "bullets":["Linux", "Networking"]},
        {"heading":"RPU", "bullets":["FreeRTOS", "Motor control"]}
      ]
    """
    columns = action.get("columns", [{}, {}])

    for i, col in enumerate(columns[:2]):
        x = Inches(0.5 + i * 6.4)
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, Inches(1.25), Inches(5.9), Inches(5.5),
        )
        fill_shape(box, DARK, MID)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.15)
        tf.clear()

        para(tf, col.get("heading", ""), 18, TEAL, True)
        for b in col.get("bullets", []):
            para(tf, f"  {b}", 13, WHITE)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# LAYOUT: key_takeaway
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def layout_key_takeaway(slide, action):
    """Single bold statement centered on the slide.

    Plan action fields:
      statement: "The APU handles application workloads; the RPU handles real-time."
      subtext: "Most Versal designs use both processors together."
    """
    statement = action.get("statement", action.get("learning_goal", ""))
    subtext = action.get("subtext", "")

    box = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.0), Inches(10.3), Inches(3.5),
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = para(tf, statement, 36, TEAL, True, PP_ALIGN.CENTER)
    p.space_after = Pt(24)
    if subtext:
        para(tf, subtext, 18, LIGHT, False, PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# LAYOUT: cards (default / legacy)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def layout_cards(slide, action):
    """2-4 info cards, stacked or side-by-side.

    Plan action fields:
      cards: [{"heading":"What", "bullets":["..."]}, ...]
    """
    cards = action.get("cards") or []
    if not cards:
        raise ValueError(
            f"Slide '{action.get('title', '<untitled>')}': layout 'cards' requires "
            f"a non-empty 'cards' array [{{\"heading\": ..., \"bullets\": [...]}}]"
        )

    if len(cards) <= 2:
        for i, card in enumerate(cards[:2]):
            x = Inches(0.75 + i * 6.0)
            _add_card_shape(slide, x, Inches(1.35), Inches(5.75), Inches(4.8), card)
    else:
        y = 1.25
        card_h = min(1.35, 5.5 / len(cards[:4]))
        for card in cards[:4]:
            _add_card_shape(slide, Inches(0.75), Inches(y), Inches(11.8), Inches(card_h), card)
            y += card_h + 0.12


def _add_card_shape(slide, x, y, w, h, card):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    fill_shape(shape, DARK, MID)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.10)
    tf.margin_bottom = Inches(0.10)
    para(tf, card.get("heading", ""), 16, GOLD, True)
    for b in card.get("bullets", []):
        para(tf, f"  {b}", 13, WHITE)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# LAYOUT: ascii_diagram
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def layout_ascii_diagram(slide, action):
    """Parse ASCII box-drawing art and render as native PPT shapes.

    Plan action fields:
      ascii_art: |
        ┌──────────┐     ┌──────────┐
        │   APU    │────►│   DDR    │
        │ Cortex-A │     │ Memory   │
        └──────────┘     └──────────┘
              │
              ▼
        ┌──────────┐
        │   RPU    │
        │ Cortex-R │
        └──────────┘

    Falls back to block_diagram layout if ascii_art is empty or parsing fails.
    """
    ascii_text = action.get("ascii_art", "")
    if not ascii_text.strip():
        layout_block_diagram(slide, action)
        return

    boxes, connectors = parse_ascii(ascii_text)
    if not boxes:
        layout_block_diagram(slide, action)
        return

    boxes_in, conns_in = boxes_to_inches(boxes, connectors)
    render_to_slide(slide, boxes_in, conns_in)

    if action.get("key_points"):
        _add_key_points_panel(slide, action["key_points"])


# ── Layout dispatcher ─────────────────────────────────────────
LAYOUTS = {
    "comparison_table": layout_comparison_table,
    "block_diagram":    layout_block_diagram,
    "ascii_diagram":    layout_ascii_diagram,
    "two_column":       layout_two_column,
    "key_takeaway":     layout_key_takeaway,
    "cards":            layout_cards,
}


# ── Main ───────────────────────────────────────────────────────
def main():
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--plan", required=True, help="Path to update_plan.json")
    ap.add_argument("--output", required=True, help="Output PPTX path")
    ap.add_argument("--strict-quality", action="store_true",
                    help="Enforce the Slide-18 quality bar: diagram/comparison_table slides "
                         "must carry source_footer/source_basis and non-empty body specs.")
    args = ap.parse_args()

    plan = safe_load_json(args.plan, "update_plan.json")
    MARP_ELIGIBLE = {"comparison_table", "two_column", "key_takeaway", "cards"}
    all_slides = [a for a in plan.get("actions", []) if a.get("type") == "add_new_slide"]
    slides = [a for a in all_slides if a.get("slide_layout") not in MARP_ELIGIBLE]

    if not slides:
        skipped = len(all_slides) - len(slides)
        if skipped:
            print(f"create_additions_deck: all {skipped} add_new_slide actions use "
                  "MARP-eligible layouts; nothing to do (use create_marp_additions.py).")
        else:
            print("create_additions_deck: no add_new_slide actions found.")
        sys.exit(0)

    if len(slides) < len(all_slides):
        print(f"create_additions_deck: processing {len(slides)} diagram-layout slides "
              f"(skipped {len(all_slides) - len(slides)} MARP-eligible slides).")

    for action in slides:
        validate_new_slide_action(action, strict_quality=args.strict_quality)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    diagram_actions = [a for a in slides if a.get("slide_layout") == "ascii_diagram" and a.get("ascii_art", "").strip()]
    if diagram_actions:
        md_path = Path(args.output).with_suffix("").as_posix() + "_diagrams.md"
        lines = ["# ASCII Diagrams for Review\n"]
        for i, action in enumerate(diagram_actions, 1):
            lines.append(f"## Diagram {i}: {action.get('title', 'Untitled')}\n")
            lines.append("```")
            lines.append(action["ascii_art"].rstrip())
            lines.append("```\n")
        Path(md_path).write_text("\n".join(lines), encoding="utf-8")
        print(f"diagrams: {md_path}")

    for action in slides:
        slide = make_slide(prs)
        add_title_bar(slide, action.get("title", "New Content"))

        layout_name = action.get("slide_layout", "cards")
        if layout_name not in LAYOUTS:
            raise ValueError(
                f"Slide '{action.get('title', '<untitled>')}': unknown slide_layout "
                f"'{layout_name}'; must be one of {sorted(LAYOUTS.keys())}"
            )
        layout_fn = LAYOUTS[layout_name]
        layout_fn(slide, action)

        # Source-attribution footer (Slide-18 quality bar element 5).
        # Drawn whenever source_footer or source_basis is set on the action,
        # regardless of strict-quality mode.
        footer_text = action.get("source_footer") or action.get("source_basis")
        if footer_text:
            _add_source_footer(slide, str(footer_text))

        add_notes(slide, action.get("speaker_notes"))

    prs.save(args.output)
    print(args.output)


def _add_source_footer(slide, text: str) -> None:
    """Add an 8-pt grey source-attribution textbox at the bottom-left."""
    box = slide.shapes.add_textbox(
        Inches(0.45), Inches(6.95), Inches(11.5), Inches(0.35),
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = para(tf, text, 8, LIGHT, False)
    p.space_after = Pt(0)


def _add_key_points_panel(slide, key_points: dict) -> None:
    """Render a right-side dark panel with heading + bullet points as visible OST.

    key_points schema:
      {"heading": "Key Specs", "bullets": ["DDR5MC: ...", "DDRMC: ...", ...]}
    """
    bullets = key_points.get("bullets", [])
    if not bullets:
        return

    heading = key_points.get("heading", "Key Points")

    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.2), Inches(1.25), Inches(3.4), Inches(5.2),
    )
    fill_shape(panel, DARK, MID)
    tf = panel.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.14)
    tf.margin_bottom = Inches(0.10)
    tf.clear()

    para(tf, heading, 14, TEAL, True)

    for bullet in bullets:
        para(tf, f"•  {bullet}", 12, WHITE)


if __name__ == "__main__":
    main()
